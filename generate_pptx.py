"""
StayHere Mobile Implementation Guide — PPTX Generator
Generates a fully styled 20-slide PowerPoint presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy, uuid

# ─── Brand Colors ───────────────────────────────────────────────────────────
BG       = RGBColor(0x08, 0x0e, 0x1e)
CARD     = RGBColor(0x11, 0x18, 0x27)
CARD2    = RGBColor(0x1a, 0x25, 0x35)
GOLD     = RGBColor(0xc9, 0xa2, 0x27)
GOLD_L   = RGBColor(0xe8, 0xb8, 0x4b)
TEAL     = RGBColor(0x0d, 0x94, 0x88)
TEAL_L   = RGBColor(0x5e, 0xea, 0xd4)
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x88, 0x99, 0xbb)
MUTED2   = RGBColor(0x55, 0x66, 0x88)
GREEN    = RGBColor(0x34, 0xd3, 0x99)
GREEN_BG = RGBColor(0x06, 0x2a, 0x1e)
ORANGE   = RGBColor(0xfb, 0x92, 0x3c)
ORANGE_BG= RGBColor(0x2a, 0x18, 0x06)
RED      = RGBColor(0xf8, 0x71, 0x71)
BLUE     = RGBColor(0x60, 0xa5, 0xfa)
BLUE_BG  = RGBColor(0x06, 0x15, 0x2a)
PURPLE   = RGBColor(0xc4, 0xb5, 0xfd)
PURPLE_BG= RGBColor(0x18, 0x10, 0x2a)
PHONE_BG = RGBColor(0x11, 0x18, 0x27)
PHONE_BR = RGBColor(0x1e, 0x2d, 0x4a)
DARK_CARD= RGBColor(0x0c, 0x14, 0x26)

# ─── Slide dimensions (widescreen 16:9) ─────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # blank

# ─── Shape helpers ───────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, radius=None):
    shape = slide.shapes.add_shape(
        1 if radius is None else 5,  # 1=rect, 5=rounded rect
        x, y, w, h
    )
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    if radius:
        # set corner rounding
        sp = shape._element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape

def add_rounded_rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0.75, radius=20000):
    shape = slide.shapes.add_shape(5, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    try:
        sp = shape._element
        prstGeom = sp.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                for child in list(avLst): avLst.remove(child)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    except: pass
    return shape

def txtbox(slide, text, x, y, w, h,
           size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
           wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    return txb

def txtbox_multi(slide, lines, x, y, w, h, default_size=12, default_color=WHITE):
    """lines: list of (text, size, bold, color, italic, align)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        text, size, bold, color, italic = line[0], line[1], line[2], line[3], line[4]
        align = line[5] if len(line) > 5 else PP_ALIGN.LEFT
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if not text:
            p.space_before = Pt(4)
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = 'Consolas' if italic else 'Segoe UI'
    return txb

def code_block(slide, lines, x, y, w, h):
    """Render a code block with dark background"""
    add_rounded_rect(slide, x, y, w, h, fill=RGBColor(0x05,0x0a,0x18),
                     line_color=RGBColor(0x1e,0x2d,0x4a), radius=12000)
    code_lines = []
    for text, color in lines:
        code_lines.append((text, 9, False, color, False))
    # use a text box with tight margin
    margin = Inches(0.15)
    txtbox_multi(slide, code_lines, x+margin, y+margin, w-margin*2, h-margin*2)

def api_row(slide, method, url, desc, x, y, w):
    """Render a color-coded API method + URL row"""
    METHOD_COLORS = {
        'GET':   (RGBColor(0x06,0x2a,0x1e), GREEN),
        'POST':  (BLUE_BG, BLUE),
        'PATCH': (ORANGE_BG, ORANGE),
        'PUT':   (PURPLE_BG, PURPLE),
        'DELETE':(RGBColor(0x2a,0x06,0x06), RED),
    }
    bg, fg = METHOD_COLORS.get(method, (CARD2, WHITE))
    badge_w = Inches(0.7)
    add_rounded_rect(slide, x, y, badge_w, Inches(0.28), fill=bg,
                     line_color=fg, line_pt=0.5, radius=8000)
    txtbox(slide, method, x+Inches(0.05), y+Inches(0.03), badge_w-Inches(0.1),
           Inches(0.22), size=9, bold=True, color=fg)
    txtbox(slide, url, x+badge_w+Inches(0.1), y+Inches(0.04), w-badge_w-Inches(0.15),
           Inches(0.22), size=9.5, color=RGBColor(0xcc,0xdd,0xff), italic=False)
    if desc:
        txtbox(slide, desc, x+badge_w+Inches(0.1), y+Inches(0.27),
               w-badge_w-Inches(0.15), Inches(0.2), size=8.5, color=MUTED)

def section_badge(slide, text, x, y):
    """Gold badge pill"""
    w = Inches(2.2)
    add_rounded_rect(slide, x, y, w, Inches(0.32), fill=RGBColor(0x1a,0x14,0x04),
                     line_color=RGBColor(0x55,0x44,0x10), radius=50000)
    txtbox(slide, text, x+Inches(0.12), y+Inches(0.04), w-Inches(0.24),
           Inches(0.24), size=9, bold=True, color=GOLD)

def info_card(slide, title, body, x, y, w, h, accent=None):
    add_rounded_rect(slide, x, y, w, h, fill=RGBColor(0x0c,0x14,0x26),
                     line_color=RGBColor(0x22,0x33,0x55), radius=10000)
    if accent:
        add_rounded_rect(slide, x, y, Inches(0.04), h,
                         fill=accent, radius=10000)
    txtbox(slide, title, x+Inches(0.15), y+Inches(0.1), w-Inches(0.25),
           Inches(0.25), size=10.5, bold=True, color=WHITE)
    txtbox(slide, body, x+Inches(0.15), y+Inches(0.36), w-Inches(0.25),
           h-Inches(0.44), size=9.5, color=MUTED, wrap=True)

def slide_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=BG)

def slide_title(slide, title, subtitle=None, x=Inches(0.7), y=Inches(1.0)):
    txtbox(slide, title, x, y, Inches(7.5), Inches(0.8),
           size=34, bold=True, color=WHITE)
    if subtitle:
        txtbox(slide, subtitle, x, y+Inches(0.82), Inches(7.2), Inches(0.5),
               size=13, color=MUTED)

# ─── Phone frame helper ───────────────────────────────────────────────────────
def phone_frame(slide, x, y, w=Inches(2.4), h=Inches(4.8)):
    """Draw a phone silhouette and return (content_x, content_y, content_w, content_h)"""
    # Outer frame
    add_rounded_rect(slide, x, y, w, h, fill=PHONE_BG,
                     line_color=PHONE_BR, line_pt=2, radius=40000)
    # Notch
    notch_w = Inches(0.7)
    notch_h = Inches(0.18)
    nx = x + (w - notch_w) / 2
    add_rounded_rect(slide, nx, y, notch_w, notch_h, fill=PHONE_BG,
                     line_color=PHONE_BR, line_pt=0.5, radius=30000)
    # Status bar stripe
    sb_h = Inches(0.22)
    add_rect(slide, x+Pt(4), y+Pt(4), w-Pt(8), sb_h, fill=RGBColor(0x05,0x09,0x14))
    txtbox(slide, '9:41', x+Inches(0.15), y+Inches(0.05), Inches(0.5), Inches(0.17),
           size=7, color=MUTED2)
    txtbox(slide, '◼◼◼ 📶', x+w-Inches(0.65), y+Inches(0.05), Inches(0.55), Inches(0.17),
           size=7, color=MUTED2)
    # Screen area
    pad = Pt(5)
    cx = x + pad
    cy = y + pad + sb_h
    cw = w - pad*2
    ch = h - pad*2 - sb_h - Inches(0.05)
    return cx, cy, cw, ch

def phone_nav_bar(slide, x, y, w, title, right_icon='🔔'):
    bar_h = Inches(0.32)
    add_rect(slide, x, y, w, bar_h, fill=RGBColor(0x0a,0x11,0x22))
    txtbox(slide, title, x+Inches(0.1), y+Inches(0.06), w-Inches(0.4), Inches(0.22),
           size=10, bold=True, color=WHITE)
    txtbox(slide, right_icon, x+w-Inches(0.3), y+Inches(0.06), Inches(0.22), Inches(0.22),
           size=10, color=MUTED)
    return y + bar_h

def phone_bottom_tabs(slide, x, y, w, active=0):
    tabs = ['🏠 Home', '🔍 Explore', '♡ Saved', '📋 Apply', '👤 Me']
    bar_h = Inches(0.42)
    add_rect(slide, x, y, w, bar_h, fill=RGBColor(0x0a,0x11,0x1e))
    tab_w = w / len(tabs)
    for i, tab in enumerate(tabs):
        color = GOLD if i == active else MUTED2
        txtbox(slide, tab, x + tab_w*i, y+Inches(0.08), tab_w, Inches(0.26),
               size=7, color=color, align=PP_ALIGN.CENTER)

def mini_listing_card(slide, x, y, w, title, price, icon='🏢', match=None):
    """Mini listing card in phone screen"""
    h = Inches(0.85)
    add_rounded_rect(slide, x, y, w, h, fill=CARD2,
                     line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    # image area
    img_w = Inches(0.65)
    add_rounded_rect(slide, x+Inches(0.06), y+Inches(0.1), img_w, Inches(0.65),
                     fill=RGBColor(0x12,0x1e,0x38), radius=5000)
    txtbox(slide, icon, x+Inches(0.06), y+Inches(0.15), img_w, Inches(0.4),
           size=18, align=PP_ALIGN.CENTER, color=WHITE)
    # text
    txtbox(slide, title, x+img_w+Inches(0.14), y+Inches(0.1), w-img_w-Inches(0.2),
           Inches(0.22), size=8, bold=True, color=WHITE)
    price_text = f'KES {price:,}/mo'
    txtbox(slide, price_text, x+img_w+Inches(0.14), y+Inches(0.34), w-img_w-Inches(0.2),
           Inches(0.18), size=8, bold=True, color=GOLD)
    if match:
        add_rounded_rect(slide, x+w-Inches(0.58), y+Inches(0.08), Inches(0.52), Inches(0.2),
                         fill=RGBColor(0x1a,0x14,0x04), line_color=GOLD, radius=20000)
        txtbox(slide, f'{match}% match', x+w-Inches(0.57), y+Inches(0.09), Inches(0.5),
               Inches(0.16), size=6.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    return y + h + Inches(0.06)

def step_pills(slide, x, y, total, active):
    """Progress step pills"""
    pill_r = Inches(0.18)
    gap = Inches(0.12)
    line_h = Inches(0.03)
    total_w = pill_r * total + gap * (total - 1) + Inches(0.4) * (total - 1)
    cx = x
    for i in range(total):
        n = i + 1
        if n < active:
            fill, tc = GREEN_BG, GREEN
        elif n == active:
            fill, tc = RGBColor(0x2a,0x1e,0x04), GOLD
        else:
            fill, tc = CARD2, MUTED2
        lc = GREEN if n < active else RGBColor(0x22,0x33,0x55)
        add_rounded_rect(slide, cx, y, pill_r, pill_r, fill=fill,
                         line_color=lc, radius=50000)
        sym = '✓' if n < active else str(n)
        col = GREEN if n < active else (GOLD if n == active else MUTED2)
        txtbox(slide, sym, cx, y+Inches(0.02), pill_r, pill_r-Inches(0.04),
               size=7.5, bold=True, color=col, align=PP_ALIGN.CENTER)
        cx += pill_r + gap
        if i < total - 1:
            lw = Inches(0.3)
            line_y = y + pill_r/2 - line_h/2
            add_rect(slide, cx, line_y, lw, line_h,
                     fill=GREEN if n < active else CARD2)
            cx += lw + gap

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────
def slide_01_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    # Gradient overlay circles (approximated with translucent rects)
    add_rounded_rect(slide, Inches(8), Inches(-1), Inches(6), Inches(6),
                     fill=RGBColor(0x0e,0x0c,0x03), radius=50000)
    add_rounded_rect(slide, Inches(-1), Inches(4), Inches(5), Inches(5),
                     fill=RGBColor(0x02,0x0e,0x0d), radius=50000)
    # Logo
    txtbox(slide, 'StayHere', Inches(3), Inches(1.0), Inches(7.3), Inches(1.5),
           size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txtbox(slide, "KENYA'S #1 RENTAL PLATFORM", Inches(3), Inches(2.3), Inches(7.3), Inches(0.4),
           size=11, color=MUTED, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Mobile App — Implementation Guide', Inches(2.5), Inches(2.85),
           Inches(8.3), Inches(0.55), size=24, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Screen-by-screen breakdown for mobile developers.\nProduction API endpoints  ·  Request/Response samples  ·  Business rules  ·  UI references.',
           Inches(3), Inches(3.5), Inches(7.3), Inches(0.6),
           size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    # Chips
    chips = ['📱 20 Screens', '🔗 25 API Endpoints', '🔐 JWT Auth (Tenant only)',
             '💳 M-Pesa STK Push', '🤖 Sage AI', '☁️ Cloudflare R2']
    chip_x = Inches(1.5)
    chip_y = Inches(4.35)
    for i, chip in enumerate(chips):
        cx = chip_x + (i % 3) * Inches(3.4)
        cy = chip_y + (i // 3) * Inches(0.42)
        add_rounded_rect(slide, cx, cy, Inches(3.1), Inches(0.32),
                         fill=CARD, line_color=RGBColor(0x22,0x33,0x55), radius=30000)
        txtbox(slide, chip, cx+Inches(0.12), cy+Inches(0.05), Inches(2.9), Inches(0.22),
               size=10, color=MUTED, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Base URL:  https://apim-dev-5c27bcf3.azure-api.net   |   Use ← → arrow keys to navigate',
           Inches(1.5), Inches(7.1), Inches(10.3), Inches(0.3),
           size=9, color=MUTED2, align=PP_ALIGN.CENTER)

def slide_02_architecture(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '🏗  ARCHITECTURE', Inches(0.7), Inches(0.3))
    slide_title(slide, 'API Services & Architecture', y=Inches(0.7))
    # Base URL box
    add_rounded_rect(slide, Inches(0.7), Inches(1.65), Inches(6.5), Inches(0.38),
                     fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, radius=8000)
    txtbox(slide, 'https://apim-dev-5c27bcf3.azure-api.net',
           Inches(0.85), Inches(1.72), Inches(6.2), Inches(0.26),
           size=11, bold=True, color=GOLD)
    # Services
    services = [
        ('Auth Service',     '/auth',        'Signup · Login · OTP verification · JWT tokens',       GOLD),
        ('Property Service', '/property',    'Listings · Bookings · Applications · Upload · Payments', TEAL),
        ('Customer Service', '/customers',   'Customer profile · Dashboard · Saved properties',       BLUE),
        ('AI Agent Service', '/aiagent',     'Sage AI — Respond & Recommend · Chat',                  PURPLE),
        ('Static Data',      '/staticdata',  'Counties · Amenities · Property types (cache these)',   RGBColor(0xf4,0x72,0xb6)),
    ]
    sy = Inches(2.15)
    for name, path, desc, color in services:
        add_rounded_rect(slide, Inches(0.7), sy, Inches(6.5), Inches(0.52),
                         fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=6000)
        add_rounded_rect(slide, Inches(0.7), sy, Inches(0.05), Inches(0.52),
                         fill=color, radius=6000)
        txtbox(slide, name, Inches(0.85), sy+Inches(0.04), Inches(2.5), Inches(0.22),
               size=10.5, bold=True, color=WHITE)
        txtbox(slide, path, Inches(3.4), sy+Inches(0.05), Inches(1.2), Inches(0.2),
               size=9.5, bold=True, color=color)
        txtbox(slide, desc, Inches(0.85), sy+Inches(0.27), Inches(6.1), Inches(0.2),
               size=8.5, color=MUTED)
        sy += Inches(0.6)
    # Auth flow box (right side)
    ax = Inches(7.6)
    add_rounded_rect(slide, ax, Inches(0.7), Inches(5.4), Inches(2.5),
                     fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=10000)
    txtbox(slide, '🔐  Auth Flow', ax+Inches(0.2), Inches(0.82), Inches(5.0), Inches(0.3),
           size=12, bold=True, color=WHITE)
    steps = [
        ('1.', 'POST /auth/signup', ' — once per user'),
        ('2.', 'POST /auth/login', ' — OTP sent via email'),
        ('3.', 'POST /auth/verify-otp', ' — receive JWT'),
        ('4.', 'Store  token, customerId, role', ''),
        ('5.', 'Check  role === "Tenant"', ' — reject others'),
        ('6.', 'Attach  Authorization: Bearer <token>', ''),
    ]
    for i, (num, key, rest) in enumerate(steps):
        sy2 = Inches(1.24) + i * Inches(0.34)
        txtbox(slide, num, ax+Inches(0.2), sy2, Inches(0.22), Inches(0.26),
               size=9.5, color=MUTED)
        txtbox(slide, key, ax+Inches(0.42), sy2, Inches(4.5), Inches(0.26),
               size=9.5, bold=False, color=GOLD_L if ('POST' in key or 'Store' in key or 'Check' in key or 'Attach' in key) else WHITE)
        if rest:
            txtbox(slide, rest, ax+Inches(0.42)+Pt(len(key)*5.5), sy2, Inches(2.5), Inches(0.26),
                   size=9.5, color=MUTED)
    # Token storage card
    add_rounded_rect(slide, ax, Inches(3.35), Inches(5.4), Inches(1.1),
                     fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=10000)
    txtbox(slide, '📦  Token Storage', ax+Inches(0.2), Inches(3.47), Inches(5.0), Inches(0.26),
           size=11, bold=True, color=WHITE)
    txtbox(slide, 'Flutter: flutter_secure_storage  |  iOS: Keychain  |  Android: Keystore\nNever store JWT in plain SharedPreferences or AsyncStorage.',
           ax+Inches(0.2), Inches(3.78), Inches(5.0), Inches(0.55),
           size=9.5, color=MUTED)
    # Header card
    add_rounded_rect(slide, ax, Inches(4.55), Inches(5.4), Inches(1.0),
                     fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=10000)
    txtbox(slide, '🌐  Every Request Must Include',
           ax+Inches(0.2), Inches(4.67), Inches(5.0), Inches(0.26),
           size=11, bold=True, color=WHITE)
    code_block(slide,
               [('Authorization: Bearer <jwt_token>', TEAL_L),
                ('Content-Type: application/json', MUTED),
                ('Accept: application/json', MUTED)],
               ax+Inches(0.1), Inches(4.98), Inches(5.2), Inches(0.52))

def slide_03_signup(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📱  SCREEN 1', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Sign Up', 'Create a new Tenant account. phoneNumber is optional. No auth header required for this request.', y=Inches(0.7))
    # API
    api_row(slide, 'POST', '/auth/signup', 'Create a new Tenant account', Inches(0.7), Inches(1.85), Inches(6.2))
    code_block(slide, [
        ('// Request body', MUTED),
        ('{', WHITE),
        ('  "fullName":    "Jane Wambui",', GOLD_L),
        ('  "email":       "jane@example.com",', GOLD_L),
        ('  "phoneNumber": "+254700000000"   // optional', MUTED),
        ('}', WHITE),
        ('', WHITE),
        ('// Response 200', MUTED),
        ('{ "message": "Account created successfully" }', GREEN),
        ('', WHITE),
        ('// Errors', MUTED),
        ('409 — Email already in use', RED),
        ('400 — Validation error (missing required field)', RED),
    ], Inches(0.7), Inches(2.3), Inches(6.2), Inches(2.7))
    info_card(slide, '💡  On Success',
              'Navigate to Login screen. No JWT is returned from signup — the user must login via OTP to receive their token.',
              Inches(0.7), Inches(5.1), Inches(6.2), Inches(0.72), accent=GOLD)
    # Phone mockup
    px, py = Inches(7.5), Inches(0.5)
    cx, cy, cw, ch = phone_frame(slide, px, py)
    # Screen content
    txtbox(slide, 'StayHere', cx, cy+Inches(0.3), cw, Inches(0.45),
           size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txtbox(slide, '👤', cx, cy+Inches(0.8), cw, Inches(0.38),
           size=22, align=PP_ALIGN.CENTER, color=TEAL_L)
    txtbox(slide, 'Create Account', cx, cy+Inches(1.22), cw, Inches(0.28),
           size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Join thousands finding their dream home', cx, cy+Inches(1.52), cw, Inches(0.2),
           size=7.5, color=MUTED, align=PP_ALIGN.CENTER)
    fields = [('FULL NAME  *', 'Jane Wambui'), ('EMAIL  *', 'jane@example.com'), ('PHONE  (optional)', '+254 700 000 000')]
    fy = cy + Inches(1.82)
    for label, placeholder in fields:
        txtbox(slide, label, cx+Inches(0.08), fy, cw-Inches(0.16), Inches(0.14),
               size=6.5, color=MUTED2, bold=True)
        add_rounded_rect(slide, cx+Inches(0.08), fy+Inches(0.16), cw-Inches(0.16), Inches(0.28),
                         fill=RGBColor(0x0c,0x14,0x26), line_color=RGBColor(0x22,0x33,0x55), radius=5000)
        txtbox(slide, placeholder, cx+Inches(0.16), fy+Inches(0.19), cw-Inches(0.32), Inches(0.22),
               size=8, color=MUTED)
        fy += Inches(0.52)
    add_rounded_rect(slide, cx+Inches(0.08), fy+Inches(0.06), cw-Inches(0.16), Inches(0.32),
                     fill=GOLD, radius=8000)
    txtbox(slide, 'Create Account  →', cx+Inches(0.08), fy+Inches(0.1), cw-Inches(0.16), Inches(0.24),
           size=9, bold=True, color=RGBColor(0x08,0x0e,0x1e), align=PP_ALIGN.CENTER)
    txtbox(slide, 'Already have an account?  Sign in', cx, fy+Inches(0.44), cw, Inches(0.2),
           size=7.5, color=MUTED, align=PP_ALIGN.CENTER)

def slide_04_login(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📱  SCREEN 2', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Login — 2-Step OTP', 'Passwordless auth. Email/phone → OTP code → JWT token.', y=Inches(0.7))
    api_row(slide, 'POST', '/auth/login', 'Step 1 — Request OTP  (field names MUST be Pascal-case)', Inches(0.7), Inches(1.82), Inches(6.2))
    code_block(slide, [
        ('// Send ONE of these — Pascal-case required', MUTED),
        ('{ "Email": "jane@example.com" }', GOLD_L),
        ('// OR', MUTED),
        ('{ "PhoneNumber": "+254700000000" }', GOLD_L),
        ('', WHITE),
        ('// Response 200', MUTED),
        ('{ "message": "OTP sent to your email" }', GREEN),
    ], Inches(0.7), Inches(2.2), Inches(6.2), Inches(1.65))
    api_row(slide, 'POST', '/auth/verify-otp', 'Step 2 — Submit 6-digit OTP → receive JWT', Inches(0.7), Inches(3.95), Inches(6.2))
    code_block(slide, [
        ('{ "Email": "jane@example.com", "Otp": "482916" }', GOLD_L),
        ('', WHITE),
        ('// Response 200', MUTED),
        ('{', WHITE),
        ('  "token":      "eyJhbGci...",', GREEN),
        ('  "role":       "Tenant",       // ⚠ reject if NOT Tenant', ORANGE),
        ('  "userId":     "usr_abc123",', TEAL_L),
        ('  "customerId": "cust_xyz789"  // save this!', TEAL_L),
        ('}', WHITE),
    ], Inches(0.7), Inches(4.3), Inches(6.2), Inches(2.3))
    # Two phones
    for idx, (step, label) in enumerate([('STEP 1', 'Enter email or phone'), ('STEP 2', 'Enter OTP code')]):
        px = Inches(7.4) + idx * Inches(2.85)
        py = Inches(0.5)
        cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(2.55), h=Inches(6.85))
        add_rounded_rect(slide, cx+Inches(0.05), cy+Inches(0.05), cw-Inches(0.1), Inches(0.22),
                         fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, radius=20000)
        txtbox(slide, step, cx, cy+Inches(0.07), cw, Inches(0.16), size=7, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txtbox(slide, 'StayHere', cx, cy+Inches(0.38), cw, Inches(0.35), size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        if idx == 0:
            txtbox(slide, 'Welcome Back 👋', cx, cy+Inches(0.84), cw, Inches(0.25), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txtbox(slide, 'Enter your email to continue', cx, cy+Inches(1.1), cw, Inches(0.2), size=8, color=MUTED, align=PP_ALIGN.CENTER)
            txtbox(slide, 'EMAIL OR PHONE', cx+Inches(0.1), cy+Inches(1.42), cw-Inches(0.2), Inches(0.16), size=6.5, color=MUTED2, bold=True)
            add_rounded_rect(slide, cx+Inches(0.1), cy+Inches(1.6), cw-Inches(0.2), Inches(0.3), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=5000)
            txtbox(slide, 'jane@example.com', cx+Inches(0.18), cy+Inches(1.65), cw-Inches(0.36), Inches(0.22), size=9, color=MUTED)
            add_rounded_rect(slide, cx+Inches(0.1), cy+Inches(2.05), cw-Inches(0.2), Inches(0.32), fill=GOLD, radius=8000)
            txtbox(slide, 'Continue  →', cx, cy+Inches(2.1), cw, Inches(0.24), size=9.5, bold=True, color=RGBColor(0x08,0x0e,0x1e), align=PP_ALIGN.CENTER)
            txtbox(slide, "New here?  Create account", cx, cy+Inches(2.52), cw, Inches(0.2), size=8, color=MUTED, align=PP_ALIGN.CENTER)
        else:
            txtbox(slide, '📨', cx, cy+Inches(0.84), cw, Inches(0.3), size=20, align=PP_ALIGN.CENTER, color=WHITE)
            txtbox(slide, 'Check your email', cx, cy+Inches(1.18), cw, Inches(0.25), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txtbox(slide, 'We sent a 6-digit code to\njane@example.com', cx, cy+Inches(1.46), cw, Inches(0.34), size=8, color=MUTED, align=PP_ALIGN.CENTER)
            # OTP boxes
            box_w = Inches(0.3)
            otp_x = cx + (cw - box_w * 6 - Inches(0.05) * 5) / 2
            for b in range(6):
                bx = otp_x + b * (box_w + Inches(0.05))
                filled = b < 3
                add_rounded_rect(slide, bx, cy+Inches(1.9), box_w, Inches(0.34),
                                 fill=DARK_CARD, line_color=GOLD if filled else RGBColor(0x22,0x33,0x55), radius=5000)
                if filled:
                    txtbox(slide, str([4,8,2][b]), bx, cy+Inches(1.93), box_w, Inches(0.28), size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            txtbox(slide, "Didn't get it?  Resend (58s)", cx, cy+Inches(2.34), cw, Inches(0.18), size=7.5, color=MUTED, align=PP_ALIGN.CENTER)
            add_rounded_rect(slide, cx+Inches(0.1), cy+Inches(2.6), cw-Inches(0.2), Inches(0.32), fill=GOLD, radius=8000)
            txtbox(slide, 'Verify Code', cx, cy+Inches(2.65), cw, Inches(0.24), size=9.5, bold=True, color=RGBColor(0x08,0x0e,0x1e), align=PP_ALIGN.CENTER)

def slide_05_home(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📱  SCREEN 3', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Home Screen', 'Load featured + available listings in parallel on mount. Images served from Cloudflare R2 CDN — no auth header for images.', y=Inches(0.7))
    api_row(slide, 'GET', '/property/listings/featured?limit=6', 'Carousel of featured/promoted properties', Inches(0.7), Inches(1.82), Inches(6.5))
    api_row(slide, 'GET', '/property/listings/available?page=1&pageSize=6', 'Recent available listings grid', Inches(0.7), Inches(2.22), Inches(6.5))
    code_block(slide, [
        ('// Listing response object (both endpoints)', MUTED),
        ('{', WHITE),
        ('  "id":              "list_abc123",', TEAL_L),
        ('  "title":           "Modern 2BR, Westlands",', GOLD_L),
        ('  "primaryImageUrl": "https://pub-89...r2.dev/...",', GOLD_L),
        ('  "images":          ["https://..."],', GOLD_L),
        ('  "price":           45000,', BLUE),
        ('  "bedrooms":        2,', BLUE),
        ('  "bathrooms":       2,', BLUE),
        ('  "location":        "Westlands, Nairobi",', GOLD_L),
        ('  "listingType":     "Rent",', GOLD_L),
        ('  "isFeatured":      true', GREEN),
        ('}', WHITE),
    ], Inches(0.7), Inches(2.65), Inches(6.5), Inches(3.15))
    info_card(slide, '💡  Image Loading',
              'Use primaryImageUrl for card thumbnail. Fall back to images[0]. R2 CDN URL — load directly as <img src>, no Bearer header needed.',
              Inches(0.7), Inches(5.88), Inches(6.5), Inches(0.72), accent=GOLD)
    # Phone
    px, py = Inches(8.0), Inches(0.45)
    cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(4.8), h=Inches(6.9))
    ny = phone_nav_bar(slide, cx, cy, cw, 'StayHere', '🔔 👤')
    # Hero
    txtbox(slide, "Kenya's #1 Platform", cx+Inches(0.1), ny+Inches(0.1), cw-Inches(0.2), Inches(0.2), size=7.5, color=MUTED)
    txtbox(slide, 'Find Your ', cx+Inches(0.1), ny+Inches(0.3), Inches(1.4), Inches(0.38), size=16, bold=True, color=WHITE)
    txtbox(slide, 'Perfect Home', cx+Inches(1.45), ny+Inches(0.3), Inches(2.0), Inches(0.38), size=16, bold=True, color=GOLD)
    # Type tabs
    tabs_y = ny + Inches(0.75)
    for i, (tab, active) in enumerate([('Rent', True), ('Buy', False), ('AirBnB', False)]):
        tw = Inches(0.75)
        tx = cx + Inches(0.1) + i * (tw + Inches(0.08))
        add_rounded_rect(slide, tx, tabs_y, tw, Inches(0.26), fill=GOLD if active else CARD2,
                         line_color=GOLD if active else RGBColor(0x22,0x33,0x55), radius=8000)
        txtbox(slide, tab, tx, tabs_y+Inches(0.04), tw, Inches(0.2),
               size=8, bold=active, color=(RGBColor(0x08,0x0e,0x1e) if active else MUTED), align=PP_ALIGN.CENTER)
    # Search bar
    sy = tabs_y + Inches(0.36)
    add_rounded_rect(slide, cx+Inches(0.1), sy, cw-Inches(0.2), Inches(0.34), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, '🔍  Search city, area, amenity…', cx+Inches(0.2), sy+Inches(0.07), cw-Inches(1.0), Inches(0.22), size=8.5, color=MUTED2)
    add_rounded_rect(slide, cx+cw-Inches(0.68), sy+Inches(0.04), Inches(0.58), Inches(0.26), fill=GOLD, radius=6000)
    txtbox(slide, 'Go', cx+cw-Inches(0.68), sy+Inches(0.06), Inches(0.58), Inches(0.22), size=9, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
    # Categories
    cats = [('🏢', 'Apt'), ('🏠', 'House'), ('🛏', 'Studio'), ('🏖', 'Villa')]
    cat_y = sy + Inches(0.45)
    cat_w = (cw - Inches(0.2) - Inches(0.06) * 3) / 4
    for i, (icon, label) in enumerate(cats):
        cx2 = cx + Inches(0.1) + i * (cat_w + Inches(0.06))
        add_rounded_rect(slide, cx2, cat_y, cat_w, Inches(0.52), fill=CARD2, line_color=RGBColor(0x22,0x33,0x55), radius=6000)
        txtbox(slide, icon, cx2, cat_y+Inches(0.04), cat_w, Inches(0.26), size=14, align=PP_ALIGN.CENTER, color=WHITE)
        txtbox(slide, label, cx2, cat_y+Inches(0.3), cat_w, Inches(0.2), size=7, color=MUTED, align=PP_ALIGN.CENTER)
    # Featured heading
    txtbox(slide, '⭐  Featured Properties', cx+Inches(0.1), cat_y+Inches(0.62), cw-Inches(0.2), Inches(0.22), size=9.5, bold=True, color=WHITE)
    card_y = cat_y + Inches(0.88)
    add_rounded_rect(slide, cx+Inches(0.1), card_y, cw-Inches(0.2), Inches(1.1), fill=CARD2, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    add_rounded_rect(slide, cx+Inches(0.1), card_y, cw-Inches(0.2), Inches(0.55), fill=RGBColor(0x10,0x1c,0x38), radius=8000)
    txtbox(slide, '🏢', cx+Inches(0.9), card_y+Inches(0.08), Inches(0.5), Inches(0.4), size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.12), card_y+Inches(0.32), Inches(0.55), Inches(0.18), fill=RGBColor(0x0a,0x20,0x30), line_color=TEAL, radius=12000)
    txtbox(slide, 'For Rent', cx+Inches(0.13), card_y+Inches(0.34), Inches(0.52), Inches(0.14), size=6.5, color=TEAL_L, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Modern 2BR, Westlands', cx+Inches(0.18), card_y+Inches(0.6), cw-Inches(0.36), Inches(0.2), size=9, bold=True, color=WHITE)
    txtbox(slide, '📍 Westlands, Nairobi', cx+Inches(0.18), card_y+Inches(0.8), cw-Inches(0.36), Inches(0.18), size=8, color=MUTED)
    txtbox(slide, 'KES 45,000/mo', cx+Inches(0.18), card_y+Inches(0.96), Inches(1.5), Inches(0.18), size=9, bold=True, color=GOLD)
    txtbox(slide, '🛏 2  🚿 2', cx+cw-Inches(0.9), card_y+Inches(0.96), Inches(0.75), Inches(0.18), size=8, color=MUTED)
    # Bottom tabs
    phone_bottom_tabs(slide, cx, cy+ch-Inches(0.45), cw, active=0)

def slide_06_explore(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📱  SCREEN 4', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Explore / Search', 'Switch to POST search when any filter is active. Infinite scroll: increment page on end-of-list.', y=Inches(0.7))
    api_row(slide, 'GET', '/property/listings/available?page=1&pageSize=10', 'Default — no filters active', Inches(0.7), Inches(1.82), Inches(6.4))
    api_row(slide, 'POST', '/property/listings/search', 'Filtered/keyword search — fire when any filter applied', Inches(0.7), Inches(2.22), Inches(6.4))
    code_block(slide, [
        ('// POST /property/listings/search — all fields optional', MUTED),
        ('{', WHITE),
        ('  "query":       "2 bedroom Kilimani",', GOLD_L),
        ('  "listingType": "Rent",  // Rent | Buy | AirBnB', GOLD_L),
        ('  "minPrice":    20000,', BLUE),
        ('  "maxPrice":    80000,', BLUE),
        ('  "bedrooms":    2,', BLUE),
        ('  "county":      "Nairobi",', GOLD_L),
        ('  "page":        1,', BLUE),
        ('  "pageSize":    10', BLUE),
        ('}', WHITE),
    ], Inches(0.7), Inches(2.65), Inches(6.4), Inches(2.75))
    api_row(slide, 'GET', '/staticdata/counties', 'County list for filter dropdown — load once, cache locally', Inches(0.7), Inches(5.48), Inches(6.4))
    info_card(slide, '📌  Infinite Scroll Logic',
              'Track currentPage in state. On scroll-to-end: increment page, append results to list. Reset to page 1 when filters change.',
              Inches(0.7), Inches(5.88), Inches(6.4), Inches(0.72), accent=TEAL)
    # Phone
    px, py = Inches(7.8), Inches(0.45)
    cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(5.0), h=Inches(6.9))
    ny = phone_nav_bar(slide, cx, cy, cw, 'Explore', '⚙️')
    # Search
    sy = ny + Inches(0.1)
    add_rounded_rect(slide, cx+Inches(0.1), sy, cw-Inches(0.86), Inches(0.34), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, '🔍  2 bedroom Kilimani…', cx+Inches(0.2), sy+Inches(0.07), cw-Inches(1.2), Inches(0.22), size=8.5, color=MUTED2)
    add_rounded_rect(slide, cx+cw-Inches(0.7), sy, Inches(0.6), Inches(0.34), fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, radius=8000)
    txtbox(slide, 'Filter', cx+cw-Inches(0.68), sy+Inches(0.07), Inches(0.56), Inches(0.22), size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Type pills
    py2 = sy + Inches(0.44)
    for i, (tab, active) in enumerate([('All', True), ('Rent', False), ('Buy', False), ('AirBnB', False)]):
        tw = Inches(0.7) if tab != 'AirBnB' else Inches(0.95)
        tx = cx + Inches(0.1) + i * (Inches(0.78))
        add_rounded_rect(slide, tx, py2, Inches(0.72), Inches(0.26), fill=GOLD if active else CARD2,
                         line_color=GOLD if active else RGBColor(0x22,0x33,0x55), radius=20000)
        txtbox(slide, tab, tx, py2+Inches(0.04), Inches(0.72), Inches(0.2),
               size=8, bold=active, color=(DARK_CARD if active else MUTED), align=PP_ALIGN.CENTER)
    txtbox(slide, '128 Properties Found', cx+Inches(0.1), py2+Inches(0.36), cw-Inches(0.2), Inches(0.22), size=9.5, bold=True, color=WHITE)
    y_card = py2 + Inches(0.64)
    for icon, title, price, beds in [('🏢', 'Modern 2BR, Westlands', 45000, '🛏2 🚿2'), ('🏠', 'Spacious 3BR, Kilimani', 70000, '🛏3 🚿2')]:
        add_rounded_rect(slide, cx+Inches(0.1), y_card, cw-Inches(0.2), Inches(0.95), fill=CARD2, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
        add_rounded_rect(slide, cx+Inches(0.1), y_card, cw-Inches(0.2), Inches(0.45), fill=RGBColor(0x10,0x1c,0x38), radius=8000)
        txtbox(slide, icon, cx+Inches(0.8), y_card+Inches(0.04), Inches(0.45), Inches(0.38), size=18, color=WHITE, align=PP_ALIGN.CENTER)
        txtbox(slide, title, cx+Inches(0.18), y_card+Inches(0.5), cw-Inches(0.36), Inches(0.22), size=9, bold=True, color=WHITE)
        txtbox(slide, f'KES {price:,}/mo', cx+Inches(0.18), y_card+Inches(0.72), Inches(1.5), Inches(0.2), size=9, bold=True, color=GOLD)
        txtbox(slide, beds, cx+cw-Inches(0.9), y_card+Inches(0.72), Inches(0.8), Inches(0.2), size=8, color=MUTED)
        y_card += Inches(1.05)
    phone_bottom_tabs(slide, cx, cy+ch-Inches(0.45), cw, active=1)

def slide_07_listing_detail(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📱  SCREEN 5', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Listing Detail', 'Entry point to rental flow. Fire view-count on mount (fire-and-forget). Load AI similar listings.', y=Inches(0.7))
    api_row(slide, 'GET', '/property/listings/{listingId}', 'Full listing data — images, amenities, terms', Inches(0.7), Inches(1.82), Inches(6.4))
    api_row(slide, 'POST', '/property/listings/{listingId}/view', 'Increment view count — fire once on mount, no body', Inches(0.7), Inches(2.22), Inches(6.4))
    api_row(slide, 'POST', '/property/listings/{listingId}/rate', 'Submit 1-5 star rating — body: { "rating": 5 }', Inches(0.7), Inches(2.62), Inches(6.4))
    api_row(slide, 'POST', '/aiagent/respondandrecommend', 'Similar properties section — body: { "query": title, "maxResults": 4 }', Inches(0.7), Inches(3.02), Inches(6.4))
    info_card(slide, '🚀  Start Rental Flow',
              '"Start Rental Process" button → navigate to Rental Flow with listingId. Require auth — redirect to Login if no token stored.',
              Inches(0.7), Inches(3.55), Inches(6.4), Inches(0.72), accent=GOLD)
    info_card(slide, '🖼  Image Gallery',
              'Use primaryImageUrl for hero. Swipe through images[] array. Add heart/save icon. Pager dots for gallery position.',
              Inches(0.7), Inches(4.35), Inches(6.4), Inches(0.72), accent=TEAL)
    info_card(slide, '⭐  Ratings',
              'Show aggregate rating from listing response. Allow users to submit a new rating. Disable submission if already rated.',
              Inches(0.7), Inches(5.15), Inches(6.4), Inches(0.72), accent=PURPLE)
    # Phone
    px, py = Inches(7.9), Inches(0.45)
    cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(5.0), h=Inches(6.9))
    # Hero image
    img_h = Inches(1.45)
    add_rect(slide, cx, cy, cw, img_h, fill=RGBColor(0x10,0x1c,0x38))
    txtbox(slide, '🏢', cx+cw/2-Inches(0.3), cy+Inches(0.35), Inches(0.6), Inches(0.7), size=28, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, '←', cx+Inches(0.1), cy+Inches(0.12), Inches(0.3), Inches(0.26), size=13, color=WHITE)
    txtbox(slide, '♡', cx+cw-Inches(0.38), cy+Inches(0.12), Inches(0.3), Inches(0.26), size=13, color=WHITE)
    txtbox(slide, '1 / 5', cx+Inches(0.1), cy+img_h-Inches(0.25), Inches(0.6), Inches(0.2), size=7.5, color=WHITE)
    # Details
    dy = cy + img_h + Inches(0.08)
    txtbox(slide, 'Modern 2BR Apartment', cx+Inches(0.1), dy, cw-Inches(0.2), Inches(0.26), size=11.5, bold=True, color=WHITE)
    txtbox(slide, '📍 Westlands, Nairobi', cx+Inches(0.1), dy+Inches(0.28), cw-Inches(0.2), Inches(0.2), size=8.5, color=MUTED)
    txtbox(slide, 'KES 45,000/mo', cx+Inches(0.1), dy+Inches(0.5), Inches(1.8), Inches(0.26), size=12, bold=True, color=GOLD)
    add_rounded_rect(slide, cx+cw-Inches(1.0), dy+Inches(0.5), Inches(0.9), Inches(0.24), fill=GREEN_BG, line_color=GREEN, radius=15000)
    txtbox(slide, 'Available', cx+cw-Inches(1.0), dy+Inches(0.52), Inches(0.9), Inches(0.2), size=7.5, color=GREEN, align=PP_ALIGN.CENTER)
    # Stats row
    sy2 = dy + Inches(0.84)
    for i, (icon, val, lbl) in enumerate([('🛏', '2', 'Bed'), ('🚿', '2', 'Bath'), ('📐', '85', 'm²')]):
        sx = cx + Inches(0.1) + i * (Inches(1.5))
        add_rounded_rect(slide, sx, sy2, Inches(1.35), Inches(0.5), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=6000)
        txtbox(slide, f'{icon}  {val} {lbl}', sx, sy2+Inches(0.12), Inches(1.35), Inches(0.26), size=9, color=WHITE, align=PP_ALIGN.CENTER)
    # Description
    txtbox(slide, 'Modern finishes, open-plan kitchen, great natural light. Secure building 24hr security…',
           cx+Inches(0.1), sy2+Inches(0.6), cw-Inches(0.2), Inches(0.34), size=8.5, color=MUTED)
    # Amenity tags
    tag_y = sy2 + Inches(1.0)
    for i, (tag, color, bg) in enumerate([('🌐 WiFi', TEAL_L, TEAL), ('🅿️ Parking', TEAL_L, TEAL), ('🛡 Security', TEAL_L, TEAL)]):
        add_rounded_rect(slide, cx+Inches(0.1)+i*Inches(1.45), tag_y, Inches(1.32), Inches(0.24), fill=RGBColor(0x06,0x1e,0x1c), line_color=TEAL, radius=15000)
        txtbox(slide, tag, cx+Inches(0.1)+i*Inches(1.45), tag_y+Inches(0.03), Inches(1.32), Inches(0.18), size=7.5, color=TEAL_L, align=PP_ALIGN.CENTER)
    # Rating
    txtbox(slide, '⭐⭐⭐⭐☆  4.2  (18 reviews)', cx+Inches(0.1), tag_y+Inches(0.34), cw-Inches(0.2), Inches(0.2), size=8.5, color=MUTED)
    # CTA
    add_rounded_rect(slide, cx+Inches(0.1), tag_y+Inches(0.6), cw-Inches(0.2), Inches(0.32), fill=GOLD, radius=8000)
    txtbox(slide, 'Start Rental Process  →', cx, tag_y+Inches(0.65), cw, Inches(0.24), size=9.5, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)

def slide_08_rental_overview(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '🏘  RENTAL FLOW', Inches(0.7), Inches(0.3))
    slide_title(slide, '6-Step Rental Process', 'Each step corresponds to API calls that advance the application status. Steps are sequential.', y=Inches(0.7))
    steps = [
        ('1', 'Book Viewing', 'POST /property/bookings', GOLD, RGBColor(0x1a,0x14,0x04)),
        ('2', 'Apply', 'POST /property/applications', TEAL_L, RGBColor(0x04,0x1a,0x18)),
        ('3', 'Upload Docs', 'POST /upload/presigned-url\nPUT to R2\nPATCH /documents', BLUE, BLUE_BG),
        ('4', 'Review Terms', 'GET /listings/{id}/terms\nPATCH /accept-terms', PURPLE, PURPLE_BG),
        ('5', 'Payment', 'POST /payment/initiate\nGET /payment/status', RGBColor(0xf4,0x72,0xb6), RGBColor(0x2a,0x06,0x18)),
        ('6', 'Welcome', 'GET /onboarding', GREEN, GREEN_BG),
    ]
    sx = Inches(0.55)
    sy = Inches(1.75)
    bw = Inches(2.05)
    bh = Inches(1.8)
    for i, (num, title, apis, color, bg) in enumerate(steps):
        bx = sx + i * (bw + Inches(0.1))
        add_rounded_rect(slide, bx, sy, bw, bh, fill=bg, line_color=color, line_pt=1.2, radius=12000)
        # Step number circle
        add_rounded_rect(slide, bx+bw/2-Inches(0.22), sy-Inches(0.22), Inches(0.44), Inches(0.44), fill=color, radius=50000)
        txtbox(slide, num, bx+bw/2-Inches(0.22), sy-Inches(0.2), Inches(0.44), Inches(0.36), size=13, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
        txtbox(slide, title, bx, sy+Inches(0.18), bw, Inches(0.32), size=12.5, bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, api_line in enumerate(apis.split('\n')):
            txtbox(slide, api_line.strip(), bx+Inches(0.08), sy+Inches(0.55)+j*Inches(0.28), bw-Inches(0.16), Inches(0.25),
                   size=8, color=MUTED, align=PP_ALIGN.CENTER)
        # Connector arrow
        if i < len(steps) - 1:
            arr_x = bx + bw + Inches(0.02)
            txtbox(slide, '→', arr_x, sy+bh/2-Inches(0.16), Inches(0.1), Inches(0.3), size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # Status flow
    sy2 = sy + bh + Inches(0.3)
    add_rounded_rect(slide, Inches(0.55), sy2, Inches(12.6), Inches(0.55), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, 'Application Status Flow:', Inches(0.72), sy2+Inches(0.09), Inches(1.6), Inches(0.26), size=9.5, bold=True, color=WHITE)
    statuses = [('DocumentsPending', ORANGE), ('→ UnderReview', MUTED), ('→ Approved', GREEN),
                ('→ TermsAccepted', GOLD), ('→ PaymentPending', ORANGE), ('→ Active ✓', GREEN)]
    sx3 = Inches(2.4)
    for text, color in statuses:
        txtbox(slide, text, sx3, sy2+Inches(0.14), Inches(1.9), Inches(0.26), size=9, bold=('→' not in text), color=color)
        sx3 += Inches(1.75)
    # Info cards row
    iy = sy2 + Inches(0.68)
    cards = [
        ('💾  Key IDs to Track', 'listingId (from listing)\nbookingId (auto from Step 1)\napplicationId (auto from Step 2)\ncustomerId (from login)', TEAL),
        ('▶️  Resume Logic', 'From My Applications, users can resume:\n• DocumentsPending → Step 3\n• Approved → Step 4\n• TermsAccepted/PaymentPending → Step 5', GOLD),
        ('⚠️  Sequential Rule', 'Steps must be completed in order. The API will reject out-of-order actions. Always pass applicationId to Steps 3–6.', ORANGE),
    ]
    cw2 = Inches(4.05)
    for i, (title, body, color) in enumerate(cards):
        info_card(slide, title, body, Inches(0.55)+i*(cw2+Inches(0.1)), iy, cw2, Inches(1.38), accent=color)

def make_step_slide(prs, screen_num, step_num, step_title, subtitle, apis_info, code_lines, tip_title, tip_body, phone_fn):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, f'🏘  STEP {step_num} OF 6', Inches(0.7), Inches(0.3))
    slide_title(slide, step_title, subtitle, y=Inches(0.7))
    y = Inches(1.82)
    for method, url, desc in apis_info:
        api_row(slide, method, url, desc, Inches(0.7), y, Inches(6.4))
        y += Inches(0.42)
    code_block(slide, code_lines, Inches(0.7), y+Inches(0.04), Inches(6.4), Inches(5.95) - y + Inches(0.7))
    if tip_title:
        info_card(slide, tip_title, tip_body, Inches(0.7), Inches(5.9), Inches(6.4), Inches(0.75), accent=GOLD)
    # Phone
    px, py = Inches(7.6), Inches(0.45)
    cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(5.3), h=Inches(6.9))
    phone_fn(slide, cx, cy, cw, ch, step_num)

def phone_step1(slide, cx, cy, cw, ch, step_num):
    ny = phone_nav_bar(slide, cx, cy, cw, 'Rental Process')
    # Step pills
    step_pills(slide, cx+Inches(0.1), ny+Inches(0.08), 6, step_num)
    py = ny + Inches(0.35)
    txtbox(slide, '📅  Schedule Viewing', cx+Inches(0.1), py, cw-Inches(0.2), Inches(0.26), size=11, bold=True, color=WHITE)
    txtbox(slide, 'Modern 2BR, Westlands · KES 45,000/mo', cx+Inches(0.1), py+Inches(0.28), cw-Inches(0.2), Inches(0.2), size=8.5, color=MUTED)
    txtbox(slide, 'VIEWING TYPE', cx+Inches(0.1), py+Inches(0.56), cw-Inches(0.2), Inches(0.16), size=7, bold=True, color=MUTED2)
    tw = (cw-Inches(0.26)) / 2
    add_rounded_rect(slide, cx+Inches(0.1), py+Inches(0.74), tw, Inches(0.34), fill=RGBColor(0x04,0x1a,0x18), line_color=TEAL, radius=8000)
    txtbox(slide, '🏠  In-Person', cx+Inches(0.1), py+Inches(0.8), tw, Inches(0.24), size=9, bold=True, color=TEAL_L, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.1)+tw+Inches(0.06), py+Inches(0.74), tw, Inches(0.34), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, '💻  Virtual', cx+Inches(0.1)+tw+Inches(0.06), py+Inches(0.8), tw, Inches(0.24), size=9, color=MUTED, align=PP_ALIGN.CENTER)
    # Mini calendar
    txtbox(slide, 'SELECT DATE', cx+Inches(0.1), py+Inches(1.18), cw, Inches(0.16), size=7, bold=True, color=MUTED2)
    add_rounded_rect(slide, cx+Inches(0.1), py+Inches(1.36), cw-Inches(0.2), Inches(0.82), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    days = 'M T W T F S S'.split()
    dates = ['16','17','18','19','20','21','22']
    dw = (cw-Inches(0.4)) / 7
    for i, (d, dt) in enumerate(zip(days, dates)):
        dx = cx+Inches(0.2)+i*dw
        txtbox(slide, d, dx, py+Inches(1.42), dw, Inches(0.18), size=6.5, color=MUTED2, align=PP_ALIGN.CENTER)
        if dt == '20':
            add_rounded_rect(slide, dx+dw*0.05, py+Inches(1.62), dw*0.9, Inches(0.22), fill=GOLD, radius=50000)
            txtbox(slide, dt, dx, py+Inches(1.64), dw, Inches(0.18), size=8, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
        else:
            txtbox(slide, dt, dx, py+Inches(1.64), dw, Inches(0.18), size=8, color=MUTED, align=PP_ALIGN.CENTER)
    # Time slots
    txtbox(slide, 'SELECT TIME', cx+Inches(0.1), py+Inches(2.28), cw, Inches(0.16), size=7, bold=True, color=MUTED2)
    times = [('10:00 AM', True), ('12:00 PM', False), ('2:00 PM', False)]
    for i, (t, active) in enumerate(times):
        tw2 = (cw-Inches(0.32)) / 3
        tx = cx+Inches(0.1)+i*(tw2+Inches(0.06))
        add_rounded_rect(slide, tx, py+Inches(2.46), tw2, Inches(0.28), fill=(RGBColor(0x1a,0x14,0x04) if active else DARK_CARD), line_color=(GOLD if active else RGBColor(0x22,0x33,0x55)), radius=6000)
        txtbox(slide, t, tx, py+Inches(2.5), tw2, Inches(0.22), size=8, bold=active, color=(GOLD if active else MUTED), align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.1), py+Inches(2.9), cw-Inches(0.2), Inches(0.32), fill=GOLD, radius=8000)
    txtbox(slide, 'Book Viewing  →', cx, py+Inches(2.95), cw, Inches(0.24), size=9.5, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)

def phone_step2(slide, cx, cy, cw, ch, step_num):
    ny = phone_nav_bar(slide, cx, cy, cw, 'Rental Process')
    step_pills(slide, cx+Inches(0.1), ny+Inches(0.08), 6, step_num)
    py = ny + Inches(0.36)
    add_rounded_rect(slide, cx+Inches(0.1), py, cw-Inches(0.2), Inches(0.62), fill=GREEN_BG, line_color=GREEN, radius=8000)
    txtbox(slide, '✅  Viewing Confirmed!', cx+Inches(0.2), py+Inches(0.08), cw-Inches(0.4), Inches(0.24), size=10, bold=True, color=GREEN)
    txtbox(slide, 'Fri Jun 20  ·  10:00 AM  ·  In-Person', cx+Inches(0.2), py+Inches(0.34), cw-Inches(0.4), Inches(0.2), size=8.5, color=MUTED)
    txtbox(slide, 'Submit Application', cx+Inches(0.1), py+Inches(0.76), cw-Inches(0.2), Inches(0.26), size=11, bold=True, color=WHITE)
    fields = [('EMPLOYMENT TYPE', 'Employed  ▾'), ('MONTHLY INCOME (KES)', '120,000'), ('EMPLOYER / BUSINESS', 'Safaricom PLC')]
    fy = py + Inches(1.1)
    for label, val in fields:
        txtbox(slide, label, cx+Inches(0.1), fy, cw-Inches(0.2), Inches(0.16), size=6.5, bold=True, color=MUTED2)
        add_rounded_rect(slide, cx+Inches(0.1), fy+Inches(0.18), cw-Inches(0.2), Inches(0.3), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=5000)
        txtbox(slide, val, cx+Inches(0.2), fy+Inches(0.22), cw-Inches(0.4), Inches(0.22), size=9, color=MUTED)
        fy += Inches(0.58)
    add_rounded_rect(slide, cx+Inches(0.1), fy+Inches(0.1), cw-Inches(0.2), Inches(0.32), fill=GOLD, radius=8000)
    txtbox(slide, 'Submit Application  →', cx, fy+Inches(0.15), cw, Inches(0.24), size=9.5, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
    txtbox(slide, '🔒  Your info is kept confidential', cx, fy+Inches(0.52), cw, Inches(0.2), size=8, color=MUTED2, align=PP_ALIGN.CENTER)

def phone_step3(slide, cx, cy, cw, ch, step_num):
    ny = phone_nav_bar(slide, cx, cy, cw, 'Rental Process')
    step_pills(slide, cx+Inches(0.1), ny+Inches(0.08), 6, step_num)
    py = ny + Inches(0.36)
    txtbox(slide, 'Upload Documents', cx+Inches(0.1), py, cw-Inches(0.2), Inches(0.26), size=11, bold=True, color=WHITE)
    txtbox(slide, 'Upload clear scans or photos of each document', cx+Inches(0.1), py+Inches(0.28), cw-Inches(0.2), Inches(0.2), size=8.5, color=MUTED)
    docs = [
        ('📄', 'ID / Passport', 'id_front.pdf  ·  1.2 MB', True),
        ('📋', 'Payslips (3 months)', 'Tap to upload PDF or image', False),
        ('🏦', 'Bank Statement', 'Last 3 months', False),
    ]
    dy = py + Inches(0.56)
    for icon, title, sub, done in docs:
        filled = GREEN_BG if done else DARK_CARD
        border = GREEN if done else RGBColor(0x22,0x33,0x44)
        add_rounded_rect(slide, cx+Inches(0.1), dy, cw-Inches(0.2), Inches(0.6), fill=filled, line_color=border, radius=8000)
        txtbox(slide, icon, cx+Inches(0.2), dy+Inches(0.12), Inches(0.4), Inches(0.38), size=18, color=WHITE)
        txtbox(slide, '✓ ' + title if done else title, cx+Inches(0.72), dy+Inches(0.1), cw-Inches(1.1), Inches(0.24), size=9.5, bold=True, color=(GREEN if done else WHITE))
        txtbox(slide, sub, cx+Inches(0.72), dy+Inches(0.35), cw-Inches(1.1), Inches(0.2), size=8, color=MUTED)
        if not done:
            add_rounded_rect(slide, cx+cw-Inches(0.72), dy+Inches(0.18), Inches(0.58), Inches(0.26), fill=GOLD, radius=6000)
            txtbox(slide, '+ Add', cx+cw-Inches(0.72), dy+Inches(0.21), Inches(0.58), Inches(0.2), size=8, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
        dy += Inches(0.68)
    # Signature area
    add_rounded_rect(slide, cx+Inches(0.1), dy, cw-Inches(0.2), Inches(0.6), fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, line_pt=0.5, radius=8000)
    txtbox(slide, '✍  Draw or upload your signature', cx+Inches(0.1), dy+Inches(0.18), cw-Inches(0.2), Inches(0.26), size=9, color=MUTED2, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.1), dy+Inches(0.68), cw-Inches(0.2), Inches(0.32), fill=MUTED2, radius=8000)
    txtbox(slide, 'Submit Documents  (1 / 4 uploaded)', cx, dy+Inches(0.73), cw, Inches(0.24), size=9, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)

def phone_step4(slide, cx, cy, cw, ch, step_num):
    ny = phone_nav_bar(slide, cx, cy, cw, 'Rental Process')
    step_pills(slide, cx+Inches(0.1), ny+Inches(0.08), 6, step_num)
    py = ny + Inches(0.36)
    add_rounded_rect(slide, cx+Inches(0.1), py, cw-Inches(0.2), Inches(0.46), fill=GREEN_BG, line_color=GREEN, radius=8000)
    txtbox(slide, '✅  Application Approved! Review your lease terms below.', cx+Inches(0.2), py+Inches(0.1), cw-Inches(0.4), Inches(0.3), size=9, bold=True, color=GREEN)
    gy = py + Inches(0.56)
    for i, (val, lbl) in enumerate([('KES 45K', 'Monthly Rent'), ('KES 90K', 'Deposit'), ('12 mo', 'Duration'), ('Jul 1', 'Move-in')]):
        gx = cx+Inches(0.1) + (i%2) * ((cw-Inches(0.28))/2 + Inches(0.06))
        if i == 2: gy = py + Inches(0.56) + Inches(0.52)
        add_rounded_rect(slide, gx, gy, (cw-Inches(0.28))/2, Inches(0.44), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=6000)
        txtbox(slide, val, gx, gy+Inches(0.02), (cw-Inches(0.28))/2, Inches(0.24), size=10.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txtbox(slide, lbl, gx, gy+Inches(0.26), (cw-Inches(0.28))/2, Inches(0.16), size=7.5, color=MUTED, align=PP_ALIGN.CENTER)
    terms_y = py + Inches(1.62)
    add_rounded_rect(slide, cx+Inches(0.1), terms_y, cw-Inches(0.2), Inches(1.0), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, 'Lease Terms', cx+Inches(0.18), terms_y+Inches(0.08), cw-Inches(0.36), Inches(0.2), size=9.5, bold=True, color=WHITE)
    txtbox(slide, 'The tenant agrees to pay rent on the 1st of each month. Late payments incur 5% penalty. Subletting requires written consent. No structural modifications…',
           cx+Inches(0.18), terms_y+Inches(0.3), cw-Inches(0.36), Inches(0.65), size=8, color=MUTED)
    accept_y = terms_y + Inches(1.1)
    add_rounded_rect(slide, cx+Inches(0.1), accept_y, cw-Inches(0.2), Inches(0.34), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=6000)
    add_rounded_rect(slide, cx+Inches(0.2), accept_y+Inches(0.07), Inches(0.22), Inches(0.22), fill=GOLD, radius=5000)
    txtbox(slide, '✓', cx+Inches(0.2), accept_y+Inches(0.07), Inches(0.22), Inches(0.22), size=9, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
    txtbox(slide, 'I have read and accept these terms', cx+Inches(0.48), accept_y+Inches(0.08), cw-Inches(0.64), Inches(0.2), size=8.5, color=MUTED)
    add_rounded_rect(slide, cx+Inches(0.1), accept_y+Inches(0.44), cw-Inches(0.2), Inches(0.32), fill=TEAL, radius=8000)
    txtbox(slide, 'Accept Terms & Continue  →', cx, accept_y+Inches(0.49), cw, Inches(0.24), size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def phone_step5(slide, cx, cy, cw, ch, step_num):
    ny = phone_nav_bar(slide, cx, cy, cw, 'Rental Process')
    step_pills(slide, cx+Inches(0.1), ny+Inches(0.08), 6, step_num)
    py = ny + Inches(0.36)
    # Summary card
    add_rounded_rect(slide, cx+Inches(0.1), py, cw-Inches(0.2), Inches(1.1), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, 'Payment Summary', cx+Inches(0.2), py+Inches(0.1), cw-Inches(0.4), Inches(0.2), size=8.5, color=MUTED)
    for i, (lbl, val) in enumerate([('First Month Rent', 'KES 45,000'), ('Security Deposit', 'KES 90,000')]):
        ry = py + Inches(0.34) + i*Inches(0.26)
        txtbox(slide, lbl, cx+Inches(0.2), ry, Inches(1.8), Inches(0.22), size=8.5, color=MUTED)
        txtbox(slide, val, cx+cw-Inches(1.2), ry, Inches(1.0), Inches(0.22), size=8.5, bold=True, color=WHITE)
    add_rect(slide, cx+Inches(0.2), py+Inches(0.84), cw-Inches(0.4), Pt(1), fill=RGBColor(0x22,0x33,0x55))
    txtbox(slide, 'Total Due', cx+Inches(0.2), py+Inches(0.9), Inches(1.5), Inches(0.22), size=10, bold=True, color=WHITE)
    txtbox(slide, 'KES 135,000', cx+cw-Inches(1.6), py+Inches(0.86), Inches(1.45), Inches(0.26), size=12, bold=True, color=GOLD)
    # M-Pesa
    txtbox(slide, 'PAYMENT METHOD', cx+Inches(0.1), py+Inches(1.2), cw-Inches(0.2), Inches(0.16), size=7, bold=True, color=MUTED2)
    add_rounded_rect(slide, cx+Inches(0.1), py+Inches(1.38), cw-Inches(0.2), Inches(0.46), fill=RGBColor(0x04,0x1a,0x18), line_color=TEAL, line_pt=1.5, radius=8000)
    txtbox(slide, '🟢  M-Pesa', cx+Inches(0.2), py+Inches(1.47), Inches(1.3), Inches(0.3), size=10, bold=True, color=TEAL_L)
    txtbox(slide, 'STK Push to your phone', cx+Inches(0.2), py+Inches(1.65), Inches(2.0), Inches(0.18), size=7.5, color=MUTED)
    txtbox(slide, '✓', cx+cw-Inches(0.4), py+Inches(1.52), Inches(0.3), Inches(0.22), size=12, bold=True, color=TEAL_L)
    # Phone input
    txtbox(slide, 'M-PESA PHONE NUMBER', cx+Inches(0.1), py+Inches(1.94), cw-Inches(0.2), Inches(0.16), size=7, bold=True, color=MUTED2)
    add_rounded_rect(slide, cx+Inches(0.1), py+Inches(2.12), cw-Inches(0.2), Inches(0.3), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=5000)
    txtbox(slide, '254 712 345 678', cx+Inches(0.2), py+Inches(2.16), cw-Inches(0.4), Inches(0.22), size=9.5, color=WHITE)
    add_rounded_rect(slide, cx+Inches(0.1), py+Inches(2.52), cw-Inches(0.2), Inches(0.32), fill=TEAL, radius=8000)
    txtbox(slide, 'Pay KES 135,000 via M-Pesa', cx, py+Inches(2.57), cw, Inches(0.24), size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, '📱  You will receive a prompt on your phone', cx, py+Inches(2.94), cw, Inches(0.2), size=8, color=MUTED2, align=PP_ALIGN.CENTER)

def phone_step6(slide, cx, cy, cw, ch, step_num):
    ny = cy  # no nav bar for welcome
    txtbox(slide, '🎉', cx, ny+Inches(0.4), cw, Inches(0.55), size=36, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Welcome Home!', cx, ny+Inches(1.0), cw, Inches(0.38), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Your rental is confirmed', cx, ny+Inches(1.4), cw, Inches(0.24), size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.1), ny+Inches(1.74), cw-Inches(0.2), Inches(1.0), fill=GREEN_BG, line_color=GREEN, radius=8000)
    txtbox(slide, '📋  Move-in Details', cx+Inches(0.2), ny+Inches(1.82), cw-Inches(0.4), Inches(0.22), size=9, bold=True, color=GREEN)
    for i, (lbl, val) in enumerate([('Address', 'Unit 4B, Westlands Apts'), ('Move-in', 'July 1, 2026'), ('Receipt', 'QJX0YK2345')]):
        ry = ny + Inches(2.08) + i*Inches(0.24)
        txtbox(slide, lbl, cx+Inches(0.2), ry, Inches(0.8), Inches(0.2), size=8, color=MUTED)
        txtbox(slide, val, cx+Inches(1.1), ry, cw-Inches(1.3), Inches(0.2), size=8, bold=True, color=(GOLD if lbl == 'Receipt' else WHITE))
    add_rounded_rect(slide, cx+Inches(0.1), ny+Inches(2.82), cw-Inches(0.2), Inches(0.7), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, '🔑  Access Instructions', cx+Inches(0.2), ny+Inches(2.9), cw-Inches(0.4), Inches(0.2), size=9, bold=True, color=WHITE)
    txtbox(slide, 'Collect keys from caretaker (James — 0720000001) on move-in day 8am–5pm.',
           cx+Inches(0.2), ny+Inches(3.12), cw-Inches(0.4), Inches(0.34), size=8, color=MUTED)
    add_rounded_rect(slide, cx+Inches(0.1), ny+Inches(3.6), cw-Inches(0.2), Inches(0.3), fill=GOLD, radius=8000)
    txtbox(slide, '📄  Download Lease', cx, ny+Inches(3.64), cw, Inches(0.24), size=9.5, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.1), ny+Inches(3.98), cw-Inches(0.2), Inches(0.3), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, '🏠  View My Dashboard', cx, ny+Inches(4.02), cw, Inches(0.24), size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

def slide_15_my_applications(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📱  SCREEN 6', Inches(0.7), Inches(0.3))
    slide_title(slide, 'My Applications', 'Track all rental applications. Resume in-progress ones. Status drives progress bar.', y=Inches(0.7))
    api_row(slide, 'GET', '/property/applications/customer/{customerId}', 'All applications for the logged-in tenant', Inches(0.7), Inches(1.82), Inches(6.4))
    code_block(slide, [
        ('// Response — array of application objects', MUTED),
        ('{', WHITE),
        ('  "applicationId": "app_ghi012",', TEAL_L),
        ('  "status":        "UnderReview",', ORANGE),
        ('  "listing": {', WHITE),
        ('    "title":          "Modern 2BR, Westlands",', GOLD_L),
        ('    "price":          45000,', BLUE),
        ('    "primaryImageUrl":"https://pub-89...r2.dev/..."', GOLD_L),
        ('  },', WHITE),
        ('  "appliedAt": "2026-06-13T09:30:00Z"', MUTED),
        ('}', WHITE),
    ], Inches(0.7), Inches(2.22), Inches(6.4), Inches(2.65))
    # Status table
    add_rounded_rect(slide, Inches(0.7), Inches(4.95), Inches(6.4), Inches(1.8), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=10000)
    txtbox(slide, 'Status → Progress Map', Inches(0.85), Inches(5.05), Inches(3.0), Inches(0.24), size=10.5, bold=True, color=WHITE)
    statuses = [('DocumentsPending', '16%', ORANGE), ('UnderReview', '33%', ORANGE),
                ('Approved', '50%', GREEN), ('TermsAccepted', '66%', GOLD),
                ('PaymentPending', '83%', ORANGE), ('Active', '100%', GREEN)]
    for i, (s, pct, c) in enumerate(statuses):
        row = i % 3; col = i // 3
        rx = Inches(0.85) + col * Inches(3.2)
        ry = Inches(5.35) + row * Inches(0.38)
        add_rounded_rect(slide, rx, ry+Inches(0.04), Inches(0.1), Inches(0.2), fill=c, radius=5000)
        txtbox(slide, s, rx+Inches(0.16), ry+Inches(0.04), Inches(1.8), Inches(0.2), size=8.5, color=WHITE)
        txtbox(slide, pct, rx+Inches(2.0), ry+Inches(0.04), Inches(0.5), Inches(0.2), size=8.5, bold=True, color=c)
    info_card(slide, '▶️  Show "Continue" Button For',
              'DocumentsPending → go to Step 3\nApproved → go to Step 4\nTermsAccepted / PaymentPending → go to Step 5',
              Inches(0.7), Inches(6.83), Inches(6.4), Inches(0.55), accent=GOLD)
    # Phone
    px, py = Inches(7.6), Inches(0.45)
    cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(5.3), h=Inches(6.9))
    ny = phone_nav_bar(slide, cx, cy, cw, 'My Applications', '🔔')
    app_y = ny + Inches(0.12)
    apps = [
        ('Modern 2BR, Westlands', 45000, 'Active', GREEN, GREEN_BG, 100, '🏢', True),
        ('Spacious 3BR, Kilimani', 70000, 'Under Review', ORANGE, ORANGE_BG, 33, '🏠', False),
        ('Beach Villa, Nyali', 120000, 'Docs Pending', GOLD, RGBColor(0x1a,0x14,0x04), 16, '🏖', True),
    ]
    for title, price, status, sc, sbg, pct, icon, resumable in apps:
        ah = Inches(1.0)
        add_rounded_rect(slide, cx+Inches(0.1), app_y, cw-Inches(0.2), ah, fill=CARD2, line_color=sc, line_pt=0.5, radius=8000)
        add_rounded_rect(slide, cx+Inches(0.2), app_y+Inches(0.1), Inches(0.55), Inches(0.55), fill=RGBColor(0x10,0x1c,0x38), radius=5000)
        txtbox(slide, icon, cx+Inches(0.2), app_y+Inches(0.1), Inches(0.55), Inches(0.55), size=18, color=WHITE, align=PP_ALIGN.CENTER)
        txtbox(slide, title, cx+Inches(0.85), app_y+Inches(0.1), cw-Inches(1.1), Inches(0.22), size=9, bold=True, color=WHITE)
        txtbox(slide, f'KES {price:,}/mo', cx+Inches(0.85), app_y+Inches(0.32), Inches(1.4), Inches(0.18), size=8, color=MUTED)
        add_rounded_rect(slide, cx+Inches(0.85), app_y+Inches(0.5), Inches(1.2), Inches(0.2), fill=sbg, line_color=sc, radius=12000)
        txtbox(slide, status, cx+Inches(0.85), app_y+Inches(0.52), Inches(1.2), Inches(0.16), size=7, bold=True, color=sc, align=PP_ALIGN.CENTER)
        # Progress bar
        bar_x = cx+Inches(0.2)
        bar_w = cw-Inches(0.4)
        add_rounded_rect(slide, bar_x, app_y+Inches(0.76), bar_w, Inches(0.1), fill=DARK_CARD, radius=5000)
        add_rounded_rect(slide, bar_x, app_y+Inches(0.76), bar_w * pct / 100, Inches(0.1), fill=sc, radius=5000)
        if resumable and status != 'Active':
            add_rounded_rect(slide, cx+cw-Inches(1.0), app_y+Inches(0.44), Inches(0.85), Inches(0.24), fill=GOLD, radius=6000)
            txtbox(slide, 'Continue →', cx+cw-Inches(1.0), app_y+Inches(0.47), Inches(0.85), Inches(0.18), size=7.5, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
        app_y += ah + Inches(0.1)
    phone_bottom_tabs(slide, cx, cy+ch-Inches(0.45), cw, active=3)

def slide_16_ai(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '🤖  SCREEN 7', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Sage AI — Two Modes', 'Mode 1: Respond & Recommend (primary). Mode 2: Chat with Sage (secondary).', y=Inches(0.7))
    # Mode 1
    add_rounded_rect(slide, Inches(0.7), Inches(1.65), Inches(6.4), Inches(2.45), fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, line_pt=1.0, radius=10000)
    txtbox(slide, '① Respond & Recommend  (Primary — first tab)', Inches(0.85), Inches(1.75), Inches(6.0), Inches(0.28), size=11, bold=True, color=GOLD)
    api_row(slide, 'POST', '/aiagent/respondandrecommend', 'AI text response + ranked listing recommendations', Inches(0.85), Inches(2.08), Inches(6.1))
    code_block(slide, [
        ('// Request', MUTED),
        ('{ "query": "2 bedroom under 100k in westlands",', GOLD_L),
        ('  "conversationId": null,  // pass returned ID on follow-ups', MUTED),
        ('  "maxResults": 5, "temperature": 0.7 }', GOLD_L),
        ('// Response  (status "000" = success)', MUTED),
        ('{ "status": "000", "message": "Nairobi vibrant...",', GREEN),
        ('  "data": { "recommendedListings": [{', WHITE),
        ('    "listingId": "494e87...", "listingCode": "L1001",', TEAL_L),
        ('    "title": "2BR apartment", "location": "Kasarani",', TEAL_L),
        ('    "price": 85000, "matchScore": 0.65,', BLUE),
        ('    "amenities": ["Parking","Balcony"] }],', GOLD_L),
        ('  "conversationId": "584119fa-...",  // store & send on next query', ORANGE),
        ('  "confidence": 0.65, "usedVectorSearch": false } }', MUTED),
    ], Inches(0.85), Inches(2.44), Inches(6.1), Inches(1.56))
    # Mode 2
    add_rounded_rect(slide, Inches(0.7), Inches(4.18), Inches(6.4), Inches(1.45), fill=RGBColor(0x04,0x14,0x18), line_color=TEAL, line_pt=1.0, radius=10000)
    txtbox(slide, '② Chat with Sage  (Secondary — second tab)', Inches(0.85), Inches(4.28), Inches(6.0), Inches(0.28), size=11, bold=True, color=TEAL_L)
    api_row(slide, 'POST', '/aiagent/chat', 'Conversational AI — maintains context via conversationHistory[]', Inches(0.85), Inches(4.6), Inches(6.1))
    code_block(slide, [
        ('{ "message": "What areas are good for families?",', GOLD_L),
        ('  "conversationHistory": [', WHITE),
        ('    { "role": "user",      "content": "..." },', MUTED),
        ('    { "role": "assistant", "content": "..." } ] }', MUTED),
        ('// Response: { "reply": "...", "listings": [...] }', GREEN),
    ], Inches(0.85), Inches(4.94), Inches(6.1), Inches(0.62))
    info_card(slide, '💡  UX: conversationId',
              'After first Respond & Recommend call, store conversationId. Pass it in all follow-up queries to maintain AI context.',
              Inches(0.7), Inches(5.7), Inches(6.4), Inches(0.65), accent=GOLD)
    # Phone
    px, py = Inches(7.6), Inches(0.45)
    cx, cy, cw, ch = phone_frame(slide, px, py, w=Inches(5.3), h=Inches(6.9))
    ny = phone_nav_bar(slide, cx, cy, cw, 'Sage AI', '● Online')
    # Mode tabs
    tab_y = ny + Inches(0.06)
    add_rounded_rect(slide, cx+Inches(0.08), tab_y, cw-Inches(0.16), Inches(0.32), fill=RGBColor(0x06,0x09,0x14), radius=8000)
    add_rounded_rect(slide, cx+Inches(0.1), tab_y+Inches(0.02), (cw-Inches(0.2))/2, Inches(0.28), fill=GOLD, radius=6000)
    txtbox(slide, '🔍  Recommend', cx+Inches(0.1), tab_y+Inches(0.05), (cw-Inches(0.2))/2, Inches(0.2), size=8.5, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
    txtbox(slide, '💬  Chat', cx+Inches(0.1)+(cw-Inches(0.2))/2, tab_y+Inches(0.07), (cw-Inches(0.2))/2, Inches(0.2), size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    # Search input
    si_y = tab_y + Inches(0.42)
    add_rounded_rect(slide, cx+Inches(0.1), si_y, cw-Inches(0.2), Inches(0.34), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, '🔍  2 bedroom under 100k in westlands', cx+Inches(0.2), si_y+Inches(0.07), cw-Inches(0.4), Inches(0.22), size=8, color=MUTED)
    # AI response bubble
    rb_y = si_y + Inches(0.44)
    add_rounded_rect(slide, cx+Inches(0.1), rb_y, cw-Inches(0.2), Inches(0.65), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    txtbox(slide, 'SAGE AI', cx+Inches(0.18), rb_y+Inches(0.07), Inches(0.8), Inches(0.16), size=6.5, bold=True, color=TEAL_L)
    txtbox(slide, 'Nairobi\'s vibrant atmosphere offers great options. Kasarani provides convenient city access. Tap L1001 to view details.',
           cx+Inches(0.18), rb_y+Inches(0.24), cw-Inches(0.36), Inches(0.38), size=8, color=MUTED)
    # Result card
    txtbox(slide, 'BEST MATCHES  ·  1 result', cx+Inches(0.1), rb_y+Inches(0.76), cw-Inches(0.2), Inches(0.18), size=7.5, bold=True, color=MUTED2)
    card_y = rb_y + Inches(0.98)
    add_rounded_rect(slide, cx+Inches(0.1), card_y, cw-Inches(0.2), Inches(0.95), fill=CARD2, line_color=RGBColor(0x22,0x33,0x55), radius=8000)
    add_rounded_rect(slide, cx+Inches(0.1), card_y, cw-Inches(0.2), Inches(0.45), fill=RGBColor(0x10,0x1c,0x38), radius=8000)
    txtbox(slide, '🏢', cx+Inches(0.8), card_y+Inches(0.04), Inches(0.45), Inches(0.38), size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+cw-Inches(1.1), card_y+Inches(0.06), Inches(0.95), Inches(0.22), fill=RGBColor(0x1a,0x14,0x04), line_color=GOLD, radius=15000)
    txtbox(slide, '65% match', cx+cw-Inches(1.1), card_y+Inches(0.08), Inches(0.95), Inches(0.18), size=7.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.12), card_y+Inches(0.26), Inches(0.5), Inches(0.18), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=5000)
    txtbox(slide, 'L1001', cx+Inches(0.12), card_y+Inches(0.28), Inches(0.5), Inches(0.14), size=6.5, color=MUTED, align=PP_ALIGN.CENTER)
    txtbox(slide, '2BR Apartment · Kasarani', cx+Inches(0.18), card_y+Inches(0.5), cw-Inches(0.36), Inches(0.22), size=9, bold=True, color=WHITE)
    txtbox(slide, 'KES 85,000/mo', cx+Inches(0.18), card_y+Inches(0.72), Inches(1.3), Inches(0.18), size=9, bold=True, color=GOLD)
    # Tags
    for i, tag in enumerate(['Parking', 'Balcony']):
        add_rounded_rect(slide, cx+cw-Inches(1.65)+i*Inches(0.78), card_y+Inches(0.72), Inches(0.7), Inches(0.2), fill=RGBColor(0x04,0x1a,0x18), line_color=TEAL, radius=10000)
        txtbox(slide, tag, cx+cw-Inches(1.65)+i*Inches(0.78), card_y+Inches(0.74), Inches(0.7), Inches(0.16), size=7, color=TEAL_L, align=PP_ALIGN.CENTER)
    # Action row
    add_rounded_rect(slide, cx+Inches(0.1), card_y+Inches(1.04), (cw-Inches(0.28))*0.65, Inches(0.3), fill=GOLD, radius=6000)
    txtbox(slide, 'View Listing  →', cx+Inches(0.1), card_y+Inches(1.07), (cw-Inches(0.28))*0.65, Inches(0.24), size=9, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, cx+Inches(0.1)+(cw-Inches(0.28))*0.65+Inches(0.08), card_y+Inches(1.04), (cw-Inches(0.28))*0.32, Inches(0.3), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=6000)
    txtbox(slide, 'Refine', cx+Inches(0.1)+(cw-Inches(0.28))*0.65+Inches(0.08), card_y+Inches(1.07), (cw-Inches(0.28))*0.32, Inches(0.24), size=9, color=MUTED, align=PP_ALIGN.CENTER)
    # Chat input bar
    chat_y = cy + ch - Inches(0.52)
    add_rect(slide, cx, chat_y, cw, Inches(0.5), fill=RGBColor(0x0a,0x11,0x1e))
    add_rounded_rect(slide, cx+Inches(0.1), chat_y+Inches(0.09), cw-Inches(0.68), Inches(0.32), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=20000)
    txtbox(slide, 'Describe what you\'re looking for…', cx+Inches(0.2), chat_y+Inches(0.14), cw-Inches(0.9), Inches(0.22), size=8, color=MUTED2)
    add_rounded_rect(slide, cx+cw-Inches(0.5), chat_y+Inches(0.09), Inches(0.36), Inches(0.32), fill=GOLD, radius=50000)
    txtbox(slide, '➤', cx+cw-Inches(0.5), chat_y+Inches(0.1), Inches(0.36), Inches(0.28), size=10, bold=True, color=DARK_CARD, align=PP_ALIGN.CENTER)

def slide_17_api_reference(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '📋  REFERENCE', Inches(0.7), Inches(0.3))
    txtbox(slide, 'Full API Reference', Inches(0.7), Inches(0.7), Inches(9.0), Inches(0.6), size=28, bold=True, color=WHITE)
    txtbox(slide, 'Base URL:', Inches(0.7), Inches(1.38), Inches(1.2), Inches(0.26), size=10, color=MUTED)
    txtbox(slide, 'https://apim-dev-5c27bcf3.azure-api.net', Inches(1.95), Inches(1.38), Inches(5.5), Inches(0.26), size=10, bold=True, color=GOLD)
    # Table header
    hx, hy = Inches(0.5), Inches(1.7)
    hw = Inches(12.5)
    add_rounded_rect(slide, hx, hy, hw, Inches(0.32), fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, radius=6000)
    for text, offset in [('Method', Inches(0.15)), ('Endpoint', Inches(1.0)), ('Screen / Purpose', Inches(5.8)), ('Auth', Inches(10.2))]:
        txtbox(slide, text, hx+offset, hy+Inches(0.07), Inches(1.8), Inches(0.2), size=9, bold=True, color=GOLD)
    rows = [
        ('POST', '/auth/signup', 'Sign Up', 'Pre-auth', ORANGE),
        ('POST', '/auth/login', 'Login Step 1', 'Pre-auth', ORANGE),
        ('POST', '/auth/verify-otp', 'Login Step 2', 'Pre-auth', ORANGE),
        ('GET',  '/property/listings/featured?limit=6', 'Home Screen', 'Bearer', GREEN),
        ('GET',  '/property/listings/available?page=1&pageSize=10', 'Home / Explore', 'Bearer', GREEN),
        ('POST', '/property/listings/search', 'Explore — filtered', 'Bearer', GREEN),
        ('GET',  '/property/listings/{listingId}', 'Listing Detail', 'Bearer', GREEN),
        ('POST', '/property/listings/{listingId}/view', 'Listing Detail — view count', 'Bearer', GREEN),
        ('POST', '/property/listings/{listingId}/rate', 'Listing Detail — rating', 'Bearer', GREEN),
        ('GET',  '/staticdata/counties', 'Explore Filters', 'Bearer', GREEN),
        ('POST', '/property/bookings', 'Rental Flow Step 1', 'Bearer', GREEN),
        ('POST', '/property/applications', 'Rental Flow Step 2', 'Bearer', GREEN),
        ('POST', '/property/upload/presigned-url', 'Rental Flow Step 3 Phase 1', 'Bearer', GREEN),
        ('PUT',  '<presigned R2 uploadUrl>', 'Rental Flow Step 3 Phase 2', 'SigV4*', TEAL_L),
        ('PATCH','/property/applications/{id}/documents', 'Rental Flow Step 3 Phase 3', 'Bearer', GREEN),
        ('GET',  '/property/listings/{listingId}/terms', 'Rental Flow Step 4', 'Bearer', GREEN),
        ('PATCH','/property/applications/{id}/accept-terms', 'Rental Flow Step 4', 'Bearer', GREEN),
        ('POST', '/property/applications/{id}/payment/initiate', 'Rental Flow Step 5 — M-Pesa', 'Bearer', GREEN),
        ('GET',  '/property/applications/{id}/payment/status', 'Rental Flow Step 5 — poll', 'Bearer', GREEN),
        ('GET',  '/property/applications/{id}/onboarding', 'Rental Flow Step 6', 'Bearer', GREEN),
        ('GET',  '/property/applications/customer/{customerId}', 'My Applications', 'Bearer', GREEN),
        ('GET',  '/customers/{customerId}', 'Customer Profile / Dashboard', 'Bearer', GREEN),
        ('POST', '/aiagent/respondandrecommend', 'AI — Respond & Recommend', 'Bearer', GREEN),
        ('POST', '/aiagent/chat', 'AI — Chat with Sage', 'Bearer', GREEN),
    ]
    METHOD_COLORS = {'GET': GREEN, 'POST': BLUE, 'PATCH': ORANGE, 'PUT': PURPLE, 'DELETE': RED}
    row_h = Inches(0.245)
    for i, (method, endpoint, screen, auth, ac) in enumerate(rows):
        ry = hy + Inches(0.32) + i * row_h
        bg = DARK_CARD if i % 2 == 0 else RGBColor(0x0a,0x11,0x1e)
        add_rect(slide, hx, ry, hw, row_h, fill=bg)
        mc = METHOD_COLORS.get(method, MUTED)
        add_rounded_rect(slide, hx+Inches(0.08), ry+Inches(0.03), Inches(0.6), Inches(0.18), fill=RGBColor(0x06,0x12,0x20), line_color=mc, radius=5000)
        txtbox(slide, method, hx+Inches(0.08), ry+Inches(0.04), Inches(0.6), Inches(0.16), size=7.5, bold=True, color=mc, align=PP_ALIGN.CENTER)
        txtbox(slide, endpoint, hx+Inches(0.82), ry+Inches(0.04), Inches(4.7), Inches(0.18), size=8, color=RGBColor(0xcc,0xdd,0xff))
        txtbox(slide, screen, hx+Inches(5.62), ry+Inches(0.04), Inches(4.4), Inches(0.18), size=8, color=MUTED)
        add_rounded_rect(slide, hx+Inches(10.12), ry+Inches(0.03), Inches(0.75), Inches(0.18), fill=RGBColor(0x06,0x12,0x20), line_color=ac, radius=5000)
        txtbox(slide, auth, hx+Inches(10.12), ry+Inches(0.04), Inches(0.75), Inches(0.16), size=7, bold=True, color=ac, align=PP_ALIGN.CENTER)
    # Footer note
    fn_y = hy + Inches(0.32) + len(rows) * row_h + Inches(0.06)
    txtbox(slide, '* SigV4 — AWS Signature v4 embedded in presigned URL. Do NOT add Authorization: Bearer header for R2 PUT uploads.',
           hx, fn_y, hw, Inches(0.22), size=8, color=MUTED2)

def slide_18_business_rules(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '⚠️  CRITICAL', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Business Rules — Must Implement', 'These rules are enforced by the API. Violating them causes errors or data corruption.', y=Inches(0.7))
    rules = [
        ('🔐', 'Tenant-Only Access', 'After verify-otp, check role === "Tenant". Block Admin, PropertyOwner, PropertyManager with an error screen. Do not store their JWT.', RED, RGBColor(0x2a,0x06,0x06)),
        ('📋', 'Sequential Rental Flow', 'Steps must be completed in order: Booking → Application → Docs → Terms → Payment. API rejects out-of-order actions.', ORANGE, ORANGE_BG),
        ('📤', 'R2 Upload Content-Type', 'The PUT to R2 MUST use the exact contentType returned in the presigned URL response. A mismatch causes 403 (SigV4 includes Content-Type in signature).', BLUE, BLUE_BG),
        ('📞', 'M-Pesa Phone Format', 'Phone must be 254XXXXXXXXX — no + prefix, no spaces. Convert: 07XXXXXXXX → 2547XXXXXXXX before submitting the payment request.', PURPLE, PURPLE_BG),
        ('🔄', 'Payment Polling — 4s Interval', 'Poll /payment/status every 4 seconds after STK Push. Stop when status is Completed or Failed. Show spinner + "Check M-Pesa on your phone".', TEAL_L, RGBColor(0x04,0x14,0x18)),
        ('🔑', 'Pascal-Case Auth Fields', 'Login and verify-otp fields MUST be Pascal-case: "Email", "PhoneNumber", "Otp". Lowercase silently fails — API returns unexpected 400.', GOLD, RGBColor(0x14,0x10,0x03)),
        ('🔁', 'Resume Applications', 'Check for existing in-progress applications before starting a new one for the same listing. Prompt user to continue or abandon existing application.', GREEN, GREEN_BG),
        ('🖼', 'Image Loading', 'R2 CDN images need no auth headers. Use primaryImageUrl, fallback to images[0]. Same URL = same file — safe to cache by URL.', RGBColor(0xf4,0x72,0xb6), RGBColor(0x2a,0x06,0x18)),
    ]
    rw = (W - Inches(1.4)) / 2
    for i, (icon, title, body, color, bg) in enumerate(rules):
        col = i % 2
        row = i // 2
        rx = Inches(0.7) + col * (rw + Inches(0.14))
        ry = Inches(1.75) + row * Inches(1.2)
        add_rounded_rect(slide, rx, ry, rw, Inches(1.1), fill=bg, line_color=color, line_pt=0.75, radius=10000)
        txtbox(slide, icon, rx+Inches(0.15), ry+Inches(0.12), Inches(0.34), Inches(0.34), size=18)
        txtbox(slide, title, rx+Inches(0.58), ry+Inches(0.12), rw-Inches(0.72), Inches(0.28), size=11, bold=True, color=color)
        txtbox(slide, body, rx+Inches(0.15), ry+Inches(0.44), rw-Inches(0.3), Inches(0.6), size=9, color=MUTED, wrap=True)

def slide_19_error_handling(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    section_badge(slide, '🛡  ERROR HANDLING', Inches(0.7), Inches(0.3))
    slide_title(slide, 'Common Errors & How to Handle', 'Standard HTTP codes. Map to user-friendly messages. Never expose raw error objects to users.', y=Inches(0.7))
    errors = [
        ('401 Unauthorized', 'Token expired or missing. Clear stored token → redirect to Login. Implement token validity check on app resume.', RED, RGBColor(0x2a,0x06,0x06)),
        ('400 Bad Request', 'Validation error — missing field or wrong format. Most common: lowercase auth field names or missing bookingId. Show the API error message body.', ORANGE, ORANGE_BG),
        ('409 Conflict', 'Duplicate record — email already registered, or application already exists for that listing. Show "You\'ve already applied" with option to view existing.', PURPLE, PURPLE_BG),
        ('403 on R2 PUT', 'Content-Type mismatch in the PUT upload. Ensure you pass the exact contentType value from the presigned URL response — not the file\'s MIME type independently.', BLUE, BLUE_BG),
        ('Payment Failed', 'M-Pesa STK timeout (30s) or insufficient funds. Show clear message. Offer "Try Again" — re-call payment/initiate. No new application needed.', GOLD, RGBColor(0x14,0x10,0x03)),
        ('5xx / Network', 'Implement exponential backoff with 3 retries for GET requests. For POST/PATCH — check idempotency before retry to avoid duplicate applications.', TEAL_L, RGBColor(0x04,0x14,0x18)),
    ]
    rw = (W - Inches(1.4)) / 2
    for i, (title, body, color, bg) in enumerate(errors):
        col = i % 2
        row = i // 2
        rx = Inches(0.7) + col * (rw + Inches(0.14))
        ry = Inches(1.75) + row * Inches(1.2)
        add_rounded_rect(slide, rx, ry, rw, Inches(1.1), fill=bg, line_color=color, line_pt=0.75, radius=10000)
        add_rounded_rect(slide, rx+Inches(0.12), ry+Inches(0.1), rw-Inches(0.24), Inches(0.28), fill=RGBColor(0x06,0x09,0x14), line_color=color, line_pt=0.5, radius=6000)
        txtbox(slide, title, rx+Inches(0.18), ry+Inches(0.14), rw-Inches(0.36), Inches(0.2), size=10, bold=True, color=color)
        txtbox(slide, body, rx+Inches(0.15), ry+Inches(0.44), rw-Inches(0.3), Inches(0.6), size=9.5, color=MUTED, wrap=True)
    # Sandbox note
    add_rounded_rect(slide, Inches(0.7), Inches(5.38), W-Inches(1.4), Inches(0.75), fill=RGBColor(0x14,0x10,0x03), line_color=GOLD, radius=8000)
    txtbox(slide, '🧪  Testing Notes', Inches(0.88), Inches(5.48), Inches(2.0), Inches(0.26), size=10.5, bold=True, color=GOLD)
    txtbox(slide, 'M-Pesa sandbox test number: +254742700000  ·  OTPs in dev are logged in Azure Function logs (not sent via email)  ·  Presigned URLs expire in 15 minutes — never cache them.',
           Inches(0.88), Inches(5.78), W-Inches(1.76), Inches(0.28), size=9.5, color=MUTED)

def slide_20_closing(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    slide_bg(slide)
    add_rounded_rect(slide, Inches(8), Inches(-1.5), Inches(7), Inches(7), fill=RGBColor(0x0e,0x0c,0x03), radius=50000)
    add_rounded_rect(slide, Inches(-1), Inches(4), Inches(6), Inches(6), fill=RGBColor(0x02,0x0e,0x0d), radius=50000)
    txtbox(slide, 'StayHere', Inches(2), Inches(1.1), Inches(9.3), Inches(1.4), size=68, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Mobile Implementation Guide', Inches(2.5), Inches(2.3), Inches(8.3), Inches(0.5), size=22, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, '🚀', Inches(5.7), Inches(2.95), Inches(1.9), Inches(0.7), size=42, align=PP_ALIGN.CENTER)
    txtbox(slide, "You're Ready to Build", Inches(2.5), Inches(3.75), Inches(8.3), Inches(0.45), size=20, color=WHITE, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Start with Auth → Home → Explore, then tackle the Rental Flow.\nSage AI and the API reference table are always here for reference.',
           Inches(2.5), Inches(4.28), Inches(8.3), Inches(0.5), size=12, color=MUTED, align=PP_ALIGN.CENTER)
    stats = [('📱', '20 Screens', 'Auth · Browse · Rental Flow\nDashboard · AI Chat'), ('🔗', '25 Endpoints', 'Across 5 Azure APIM\nservice routes'), ('🎯', 'Production Ready', 'Real URLs · Real payloads\nReal business rules')]
    sw = Inches(3.8)
    for i, (icon, title, desc) in enumerate(stats):
        sx = Inches(0.7) + i * (sw + Inches(0.15))
        add_rounded_rect(slide, sx, Inches(5.05), sw, Inches(1.2), fill=DARK_CARD, line_color=RGBColor(0x22,0x33,0x55), radius=10000)
        txtbox(slide, icon, sx, Inches(5.15), sw, Inches(0.3), size=18, align=PP_ALIGN.CENTER)
        txtbox(slide, title, sx, Inches(5.47), sw, Inches(0.26), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txtbox(slide, desc, sx, Inches(5.74), sw, Inches(0.44), size=9, color=MUTED, align=PP_ALIGN.CENTER)
    txtbox(slide, 'Base URL: https://apim-dev-5c27bcf3.azure-api.net   ·   R2 CDN: pub-89021444cd6a4a81bacfa5c1dd0427ea.r2.dev',
           Inches(1.0), Inches(7.22), Inches(11.3), Inches(0.24), size=9, color=MUTED2, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# BUILD
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = new_prs()
    print('Building slides...')
    slide_01_cover(prs);            print('  1/20 Cover')
    slide_02_architecture(prs);     print('  2/20 Architecture')
    slide_03_signup(prs);           print('  3/20 Sign Up')
    slide_04_login(prs);            print('  4/20 Login')
    slide_05_home(prs);             print('  5/20 Home Screen')
    slide_06_explore(prs);          print('  6/20 Explore')
    slide_07_listing_detail(prs);   print('  7/20 Listing Detail')
    slide_08_rental_overview(prs);  print('  8/20 Rental Flow Overview')

    # Steps 1-6 reuse make_step_slide + custom phone functions
    make_step_slide(prs, 9, 1, 'Step 1 — Book a Viewing',
        'User picks date, time and viewing type. Creates a viewing booking linked to the listing.',
        [('POST', '/property/bookings', 'Create a viewing appointment — requires Bearer auth')],
        [('// Request', MUTED), ('{', WHITE),
         ('  "listingId":   "{{listingId}}",', GOLD_L),
         ('  "customerId":  "{{customerId}}",', GOLD_L),
         ('  "viewingDate": "2026-07-01",', GOLD_L),
         ('  "viewingTime": "10:00",', GOLD_L),
         ('  "viewingType": "InPerson"  // InPerson | Virtual', MUTED),
         ('}', WHITE),('', WHITE),
         ('// Response 201', MUTED),
         ('{', WHITE),
         ('  "bookingId":   "bk_def456",', TEAL_L),
         ('  "status":      "Confirmed",', GREEN),
         ('  "viewingDate": "2026-07-01T10:00:00"', GOLD_L),
         ('}', WHITE)],
        '💾  After Success', 'Store bookingId in state. Advance to Step 2 (Apply). Show confirmation card with date/time.', phone_step1)
    print('  9/20 Step 1 Book Viewing')

    make_step_slide(prs, 10, 2, 'Step 2 — Submit Application',
        'Tenant formally applies with employment and income details. Requires bookingId from Step 1.',
        [('POST', '/property/applications', 'Submit a formal rental application — requires Bearer auth')],
        [('// Request', MUTED), ('{', WHITE),
         ('  "listingId":      "{{listingId}}",', GOLD_L),
         ('  "customerId":     "{{customerId}}",', GOLD_L),
         ('  "bookingId":      "{{bookingId}}",', GOLD_L),
         ('  "employmentType": "Employed",  // Employed|SelfEmployed|Unemployed|Student', MUTED),
         ('  "monthlyIncome":  120000,', BLUE),
         ('  "employer":       "Safaricom PLC"', GOLD_L),
         ('}', WHITE),('', WHITE),
         ('// Response 201', MUTED),
         ('{', WHITE),
         ('  "applicationId": "app_ghi012",', TEAL_L),
         ('  "status":        "DocumentsPending"', ORANGE),
         ('}', WHITE)],
        '💾  Store applicationId', 'applicationId is required for ALL subsequent steps (3-6). Show success animation before advancing.', phone_step2)
    print(' 10/20 Step 2 Apply')

    make_step_slide(prs, 11, 3, 'Step 3 — Upload Documents',
        'Two-phase upload: presigned URL from API, then PUT directly to Cloudflare R2. Never POST files to the API.',
        [('POST', '/property/upload/presigned-url', 'Phase 1: get time-limited R2 upload URL — Bearer required'),
         ('PUT', '<presigned uploadUrl>', 'Phase 2: PUT raw file bytes to R2 — SigV4 auth embedded in URL, no Bearer'),
         ('PATCH', '/property/applications/{id}/documents', 'Phase 3: submit collected publicUrls to application')],
        [('// Phase 1 Request — POST /property/upload/presigned-url', MUTED),
         ('{ "folder": "applications/docs", "fileName": "id.pdf", "contentType": "application/pdf" }', GOLD_L),
         ('// Phase 1 Response', MUTED),
         ('{ "uploadUrl":  "https://stayhere...r2.cloudflarestorage.com/...?X-Amz-Signature=...",', TEAL_L),
         ('  "publicUrl":  "https://pub-89021444cd6a4a81bacfa5c1dd0427ea.r2.dev/applications/docs/id.pdf",', TEAL_L),
         ('  "contentType": "application/pdf" }', MUTED),('', WHITE),
         ('// Phase 2 — PUT to uploadUrl with EXACT Content-Type header', MUTED),
         ('// ⚠  No Bearer token! SigV4 is embedded in uploadUrl', ORANGE),('', WHITE),
         ('// Phase 3 Request — PATCH /applications/{id}/documents', MUTED),
         ('{ "idDocumentUrl": "https://pub-89...r2.dev/...",', GOLD_L),
         ('  "payslipUrls":   ["...", "...", "..."],', GOLD_L),
         ('  "bankStatementUrl": "...", "signatureUrl": "..." }', GOLD_L)],
        '⚠️  Content-Type Must Match', 'Use EXACTLY the contentType from the presigned URL response when doing the PUT — not the file\'s MIME type independently. Mismatch = 403.', phone_step3)
    print(' 11/20 Step 3 Upload Docs')

    make_step_slide(prs, 12, 4, 'Step 4 — Review & Accept Terms',
        'Only shown when applicationStatus === "Approved". Display lease terms, get acceptance.',
        [('GET', '/property/listings/{listingId}/terms', 'Fetch lease terms for display — Bearer required'),
         ('PATCH', '/property/applications/{id}/accept-terms', 'User accepts terms — no request body — Bearer required')],
        [('// GET /listings/{id}/terms response', MUTED),
         ('{', WHITE),
         ('  "termsId":             "terms_abc",', TEAL_L),
         ('  "rentAmount":          45000,', BLUE),
         ('  "depositAmount":       90000,', BLUE),
         ('  "leaseDurationMonths": 12,', BLUE),
         ('  "moveInDate":          "2026-07-01",', GOLD_L),
         ('  "termsContent":        "The tenant agrees to..."', GOLD_L),
         ('}', WHITE),('', WHITE),
         ('// PATCH accept-terms — no body needed', MUTED),
         ('// Response 200', MUTED),
         ('{ "status": "TermsAccepted" }', GREEN),
         ('// → Advance to Step 5 (Payment)', MUTED)],
        '⚠️  Status Guard', 'Only show Step 4 if applicationStatus === "Approved". If "UnderReview" → show waiting screen. If "Rejected" → show rejection message.', phone_step4)
    print(' 12/20 Step 4 Terms')

    make_step_slide(prs, 13, 5, 'Step 5 — M-Pesa Payment',
        'STK Push to tenant phone. Poll every 4 seconds until Completed or Failed.',
        [('POST', '/property/applications/{id}/payment/initiate', 'Trigger M-Pesa STK Push — Bearer required'),
         ('GET', '/property/applications/{id}/payment/status', 'Poll every 4 seconds — stop when status !== "Pending"'),
         ('GET', '/property/applications/{id}/onboarding', 'After payment Completed — fetch Step 6 welcome data')],
        [('// POST initiate — Request', MUTED),
         ('{ "phoneNumber": "254712345678",  // no + prefix', GOLD_L),
         ('  "amount": 45000 }', GOLD_L),
         ('// Response', MUTED),
         ('{ "checkoutRequestId": "ws_CO_...", "message": "STK Push sent" }', GREEN),('', WHITE),
         ('// GET /payment/status — poll every 4 seconds', MUTED),
         ('{ "status": "Pending" | "Completed" | "Failed",', ORANGE),
         ('  "mpesaReceiptNumber": "QJX0YK2345" }', TEAL_L),('', WHITE),
         ('// Stop polling when status !== "Pending"', MUTED),
         ('// On Completed → call /onboarding → advance to Step 6', GREEN),
         ('// On Failed → show retry button → re-call initiate', RED)],
        '📱  Phone Format', 'Convert 07XXXXXXXX → 2547XXXXXXXX before submitting. No + prefix. No spaces. Timeout after 30s — user did not approve prompt.', phone_step5)
    print(' 13/20 Step 5 Payment')

    make_step_slide(prs, 14, 6, 'Step 6 — Welcome Home!',
        'Payment confirmed. Fetch onboarding info and celebrate — application is now Active.',
        [('GET', '/property/applications/{id}/onboarding', 'Fetch move-in instructions after payment confirmed')],
        [('// Response', MUTED),
         ('{', WHITE),
         ('  "applicationId":    "app_ghi012",', TEAL_L),
         ('  "status":           "Active",', GREEN),
         ('  "moveInDate":       "2026-07-01",', GOLD_L),
         ('  "propertyAddress":  "Unit 4B, Westlands Apts",', GOLD_L),
         ('  "accessInstructions": "Collect keys from caretaker...",', GOLD_L),
         ('  "emergencyContact": "+254720000001",', GOLD_L),
         ('  "leaseDocumentUrl": "https://pub-89...r2.dev/..."', TEAL_L),
         ('}', WHITE),('', WHITE),
         ('// applicationStatus is now "Active"', GREEN),
         ('// Clear rental flow state after this screen', MUTED),
         ('// Offer "Download Lease" using leaseDocumentUrl', MUTED)],
        '🎉  After This Step', 'Navigate user to My Applications or Dashboard. Clear all rental flow state (listingId, bookingId, applicationId). Offer push notification opt-in.', phone_step6)
    print(' 14/20 Step 6 Welcome')

    slide_15_my_applications(prs); print(' 15/20 My Applications')
    slide_16_ai(prs);               print(' 16/20 AI Chat')
    slide_17_api_reference(prs);    print(' 17/20 API Reference')
    slide_18_business_rules(prs);   print(' 18/20 Business Rules')
    slide_19_error_handling(prs);   print(' 19/20 Error Handling')
    slide_20_closing(prs);          print(' 20/20 Closing')

    out = r'd:\Projects\StayHereProject_MVP\stayhere-mvp\StayHere_Mobile_Implementation_Guide.pptx'
    prs.save(out)
    print(f'\n✅  Saved: {out}')

if __name__ == '__main__':
    build()
