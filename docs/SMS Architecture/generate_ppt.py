from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)   # dark navy – slide backgrounds / header bars
TEAL       = RGBColor(0x00, 0x8B, 0x9A)   # accent / highlights
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xFA)   # very light teal for card backgrounds
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x2E, 0x2E, 0x2E)
MID_GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
GREEN      = RGBColor(0x00, 0x96, 0x57)
AMBER      = RGBColor(0xFF, 0x8F, 0x00)
RED        = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def fill_slide_bg(slide, color):
    """Fill a slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=12, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_label_value(slide, label, value, left, top, width=3.5, lbl_size=9, val_size=11):
    add_textbox(slide, label, left, top, width, 0.25,
                font_size=lbl_size, bold=False, color=TEAL)
    add_textbox(slide, value, left, top + 0.22, width, 0.35,
                font_size=val_size, bold=True, color=WHITE)


def header_bar(slide, title, subtitle=""):
    """Dark navy header bar at top of slide."""
    add_rect(slide, 0, 0, 13.33, 1.3, NAVY)
    add_textbox(slide, title, 0.35, 0.12, 10, 0.6,
                font_size=24, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.35, 0.68, 10, 0.45,
                    font_size=13, bold=False, color=LIGHT_TEAL, italic=True)
    # Teal accent line under header
    add_rect(slide, 0, 1.3, 13.33, 0.05, TEAL)


def bullet_block(slide, title, bullets, left, top, width, height,
                 bg_color=LIGHT_GRAY, title_color=NAVY, bullet_color=DARK_GRAY,
                 title_size=11, bullet_size=10, corner_teal=True):
    """Card with title + bullet list."""
    add_rect(slide, left, top, width, height, bg_color)
    if corner_teal:
        add_rect(slide, left, top, 0.05, height, TEAL)
    add_textbox(slide, title, left + 0.12, top + 0.08, width - 0.2, 0.3,
                font_size=title_size, bold=True, color=title_color)
    y = top + 0.38
    for b in bullets:
        add_textbox(slide, f"• {b}", left + 0.12, y, width - 0.2, 0.25,
                    font_size=bullet_size, color=bullet_color)
        y += 0.27


def icon_card(slide, icon, label, desc, left, top, w=2.5, h=1.4):
    add_rect(slide, left, top, w, h, LIGHT_GRAY, TEAL, 1)
    add_textbox(slide, icon, left + 0.1, top + 0.08, 0.5, 0.5,
                font_size=20, color=TEAL)
    add_textbox(slide, label, left + 0.6, top + 0.1, w - 0.7, 0.3,
                font_size=11, bold=True, color=NAVY)
    add_textbox(slide, desc, left + 0.1, top + 0.55, w - 0.2, 0.75,
                font_size=9, color=MID_GRAY)


def flow_step(slide, number, label, left, top, w=1.9, h=0.9):
    add_rect(slide, left, top, w, h, NAVY)
    add_textbox(slide, number, left + 0.08, top + 0.06, 0.3, 0.3,
                font_size=14, bold=True, color=TEAL)
    add_textbox(slide, label, left + 0.38, top + 0.08, w - 0.45, 0.75,
                font_size=9, bold=False, color=WHITE, wrap=True)


def arrow_h(slide, left, top, length=0.35):
    """Simple right-pointing arrow as a thin teal rect + triangle-ish text."""
    add_rect(slide, left, top + 0.08, length, 0.04, TEAL)
    add_textbox(slide, "▶", left + length - 0.05, top - 0.01, 0.25, 0.25,
                font_size=8, color=TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE / COVER
# ══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide1, NAVY)

# Full-width teal accent stripe
add_rect(slide1, 0, 5.8, 13.33, 0.12, TEAL)

# Logo placeholder (text-based)
add_textbox(slide1, "ELMA", 0.5, 0.4, 3, 0.7,
            font_size=36, bold=True, color=TEAL)
add_textbox(slide1, "by StayHere", 0.5, 1.0, 3, 0.35,
            font_size=13, bold=False, color=LIGHT_TEAL)

# Main title
add_textbox(slide1, "SMS Platform Architecture", 0.5, 1.9, 12, 1.0,
            font_size=40, bold=True, color=WHITE)
add_textbox(slide1,
            "Multi-Channel · Multi-Country · High-Reliability SMS Delivery",
            0.5, 2.85, 12, 0.5,
            font_size=17, bold=False, color=LIGHT_TEAL, italic=True)

# Sub-line
add_textbox(slide1,
            "Board Presentation  |  Architecture & Engineering Overview",
            0.5, 3.55, 12, 0.4,
            font_size=13, color=RGBColor(0xAA, 0xCC, 0xDD))

# Date
add_textbox(slide1, "June 2026", 0.5, 6.1, 4, 0.35,
            font_size=11, color=LIGHT_TEAL)

# Decorative right block
add_rect(slide1, 10.5, 0, 2.83, 7.5, RGBColor(0x05, 0x28, 0x55))
add_textbox(slide1, "SMS\n◆\nINFRA", 10.7, 2.8, 2.5, 2,
            font_size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide2, WHITE)
header_bar(slide2, "Executive Summary",
           "What the Elma SMS Platform does and why it matters")

# 3 value-prop cards
cards = [
    ("◆ Multi-Channel Reach",
     ["App, USSD, and Web channels", "Unified pipeline for all entry points",
      "No channel-specific code paths"]),
    ("◆ Queue-Isolated Routing",
     ["One queue per Country × Bank × Category", "Independent scaling per segment",
      "Fault containment across regions"]),
    ("◆ Reliable Delivery",
     ["Persistent RabbitMQ queues (no data loss)", "Dead-letter queues for failures",
      "Exponential backoff retry logic"]),
]
x_positions = [0.35, 4.55, 8.75]
for (title, bullets), x in zip(cards, x_positions):
    bullet_block(slide2, title, bullets,
                 left=x, top=1.55, width=3.85, height=2.1,
                 bg_color=LIGHT_GRAY, title_color=NAVY)

# Key metrics row
add_rect(slide2, 0.35, 3.9, 12.6, 1.6, NAVY)
metrics = [
    ("3", "Entry Channels"),
    ("N×N×N", "Isolated Queues"),
    ("2", "Delivery Protocols\n(SMPP + SDP)"),
    ("3", "Deployment Tiers"),
    ("0", "Message Loss\n(Dead-letter safety)"),
]
for i, (val, lbl) in enumerate(metrics):
    x = 0.6 + i * 2.5
    add_textbox(slide2, val, x, 4.05, 2, 0.55,
                font_size=32, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(slide2, lbl, x, 4.55, 2, 0.7,
                font_size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

# Footer note
add_textbox(slide2,
            "Designed for multi-country banking operations — scalable, resilient, protocol-agnostic",
            0.35, 5.7, 12.6, 0.4,
            font_size=10.5, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – SYSTEM COMPONENTS
# ══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide3, WHITE)
header_bar(slide3, "System Components",
           "11 components across three deployment tiers")

components = [
    ("App Channel",       "Native mobile",   "Mobile client transactions"),
    ("USSD Channel",      "Telecom gateway",  "Carrier/hosted USSD sessions"),
    ("Web Channel",       "Browser/Cloud",    "Direct SMS API calls"),
    ("Elma Core",         "C# .NET / IIS",    "Central transaction engine"),
    ("SMS API",           "C# .NET / Docker", "Web channel entry point"),
    ("RabbitMQ Wrapper",  "C# .NET / Docker", "Queue name resolution & publish"),
    ("RabbitMQ Broker",   "Docker/VM",        "Durable message persistence"),
    ("Queue Processor×N", "C# Worker / Docker","One per queue — AMQP consumer"),
    ("SMS Splitter GW",   "C# .NET / IIS",    "MNO routing by MSISDN prefix"),
    ("MNO Passthrough",   "Node.js / PM2",    "HTTP-to-SMPP/SDP adapter per MNO"),
    ("MNO SMSC",          "Node.js / PM2",    "SMPP/SDP client delivery module"),
]

col_w = 3.85
row_h = 0.72
cols = 3
for idx, (name, runtime, role) in enumerate(components):
    col = idx % cols
    row = idx // cols
    x = 0.35 + col * (col_w + 0.3)
    y = 1.55 + row * (row_h + 0.15)
    bg = LIGHT_TEAL if row == 0 else LIGHT_GRAY
    add_rect(slide3, x, y, col_w, row_h, bg, TEAL, 0.5)
    add_rect(slide3, x, y, 0.05, row_h, TEAL)
    add_textbox(slide3, name, x + 0.12, y + 0.04, col_w - 0.2, 0.28,
                font_size=10.5, bold=True, color=NAVY)
    add_textbox(slide3, runtime, x + 0.12, y + 0.28, col_w - 0.2, 0.2,
                font_size=8.5, color=TEAL, italic=True)
    add_textbox(slide3, role, x + 0.12, y + 0.44, col_w - 0.2, 0.22,
                font_size=8.5, color=MID_GRAY)

# Legend
add_textbox(slide3, "Light teal = Channel / Entry  |  Light grey = Processing / Delivery",
            0.35, 6.9, 12.6, 0.35, font_size=9, color=MID_GRAY, italic=True,
            align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – THREE ENTRY CHANNELS
# ══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide4, WHITE)
header_bar(slide4, "Three Entry Channels",
           "All converge on the same asynchronous pipeline")

channels = [
    ("📱 App Channel", "POST /api/transactions",
     ["Sent to Elma Core",
      "Transaction validated & SMS need detected",
      "Enqueued asynchronously",
      "202 Accepted returned immediately",
      "Typical triggers: OTP, balance alert, confirmation"]),
    ("📟 USSD Channel", "POST /api/ussd/session",
     ["Mirrors App channel exactly",
      "USSD response returned immediately",
      "SMS delivery is fully asynchronous",
      "No blocking of USSD session",
      "Same queue pipeline as App"]),
    ("🌐 Web Channel", "POST /api/sms/send",
     ["Hits SMS API directly (bypasses Elma Core)",
      "Pre-processed / external sender use-case",
      "SMS API → RabbitMQ Wrapper → Queue",
      "Independently scalable endpoint",
      "Suitable for bulk or direct API callers"]),
]

for i, (title, endpoint, bullets) in enumerate(channels):
    x = 0.35 + i * 4.35
    # Channel header card
    add_rect(slide4, x, 1.55, 4.1, 0.65, NAVY)
    add_textbox(slide4, title, x + 0.15, 1.6, 3.8, 0.3,
                font_size=13, bold=True, color=WHITE)
    add_textbox(slide4, endpoint, x + 0.15, 1.88, 3.8, 0.25,
                font_size=9, color=TEAL, italic=True)
    # Bullet card
    add_rect(slide4, x, 2.25, 4.1, 3.2, LIGHT_GRAY)
    add_rect(slide4, x, 2.25, 0.05, 3.2, TEAL)
    y = 2.35
    for b in bullets:
        add_textbox(slide4, f"• {b}", x + 0.12, y, 3.85, 0.28,
                    font_size=9.5, color=DARK_GRAY)
        y += 0.55

# Convergence note
add_rect(slide4, 0.35, 5.6, 12.6, 0.75, NAVY)
add_textbox(slide4,
            "All three channels converge on the same RabbitMQ queue pipeline — "
            "ensuring identical delivery guarantees regardless of entry point.",
            0.5, 5.72, 12.3, 0.5,
            font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – END-TO-END MESSAGE FLOW
# ══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide5, WHITE)
header_bar(slide5, "End-to-End Message Flow",
           "App & USSD path — 8 steps from transaction to delivery")

steps = [
    ("1", "Transaction\nPosted"),
    ("2", "SMS Output\nDetected"),
    ("3", "Queue\nName Resolved"),
    ("4", "202 Accepted\nReturned"),
    ("5", "Queue\nProcessor\nConsumes"),
    ("6", "SMS Splitter\nRoutes by\nMSISDN"),
    ("7", "MNO\nPassthrough\nDelivers"),
    ("8", "ACK Received\n& Msg Removed"),
]

step_w = 1.42
step_h = 1.0
y_steps = 1.7
x_start = 0.35
gap = 0.08

for i, (num, label) in enumerate(steps):
    x = x_start + i * (step_w + gap + 0.35)
    flow_step(slide5, num, label, x, y_steps, step_w, step_h)
    if i < len(steps) - 1:
        arrow_h(slide5, x + step_w, y_steps + 0.45, 0.35)

# Details row
details = [
    ("Elma Core receives POST /api/transactions",       0.35),
    ("Inspects output — OTP, alert, confirmation",      1.9),
    ("Queue: {COUNTRY}_{BANK}_{CATEGORY}_SMS",          3.45),
    ("HTTP client unblocked immediately",               4.95),
    ("Long-lived AMQP consumer subscription",          6.45),
    ("Prefix match: +25677 → Airtel Uganda",            7.95),
    ("SMPP PDU or SDP HTTP to MNO SMSC",               9.5),
    ("delivery_mode=2 persistence ensured",            10.95),
]

add_rect(slide5, 0.35, 2.88, 12.6, 1.5, LIGHT_GRAY)
for text, xpos in details:
    add_textbox(slide5, text, xpos, 2.96, 1.5, 1.3,
                font_size=7.8, color=MID_GRAY, wrap=True)

# Web channel note
add_rect(slide5, 0.35, 4.55, 12.6, 0.9, LIGHT_TEAL)
add_rect(slide5, 0.35, 4.55, 0.05, 0.9, TEAL)
add_textbox(slide5, "Web Channel Path", 0.5, 4.62, 3, 0.28,
            font_size=10, bold=True, color=NAVY)
add_textbox(slide5,
            "POST /api/sms/send  →  SMS API  →  RabbitMQ API Wrapper  →  Queue  →  "
            "Steps 5–8 identical to App/USSD",
            0.5, 4.88, 12.2, 0.45,
            font_size=10, color=DARK_GRAY)

# Key insight
add_rect(slide5, 0.35, 5.62, 12.6, 0.72, NAVY)
add_textbox(slide5,
            "Async-first design: MNO back-pressure never reaches the transaction layer — "
            "Elma Core returns 202 before any SMS delivery attempt.",
            0.5, 5.75, 12.2, 0.45,
            font_size=10.5, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – QUEUE NAMING & ROUTING
# ══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide6, WHITE)
header_bar(slide6, "Queue Naming & Routing Strategy",
           "Granular isolation enables independent scaling, tuning, and fault containment")

# Format box
add_rect(slide6, 0.35, 1.55, 12.6, 0.9, NAVY)
add_textbox(slide6, "Queue Name Format:", 0.55, 1.62, 3, 0.3,
            font_size=11, bold=True, color=TEAL)
add_textbox(slide6, "{COUNTRY}  _  {BANK_ID}  _  {CATEGORY}  _  SMS",
            0.55, 1.88, 12, 0.45,
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Examples
examples = [
    ("UGANDA_99_OTP_SMS",      "Uganda · Bank 99 · One-Time Password"),
    ("KENYA_12_PROMO_SMS",     "Kenya · Bank 12 · Promotional messages"),
    ("TANZANIA_07_ALERT_SMS",  "Tanzania · Bank 7 · Account alerts"),
]
add_textbox(slide6, "Examples", 0.35, 2.58, 3, 0.3,
            font_size=11, bold=True, color=NAVY)
for i, (qname, desc) in enumerate(examples):
    y = 2.88 + i * 0.45
    add_rect(slide6, 0.35, y, 4.5, 0.38, LIGHT_TEAL)
    add_textbox(slide6, qname, 0.45, y + 0.05, 4.3, 0.28,
                font_size=10.5, bold=True, color=NAVY)
    add_textbox(slide6, desc, 4.95, y + 0.07, 5, 0.28,
                font_size=10, color=MID_GRAY, italic=True)

# Benefits
benefits = [
    ("Independent Scaling",
     "Per-segment throughput without cross-contamination"),
    ("Priority Tuning",
     "OTP: high prefetch · PROMO: low prefetch"),
    ("Fault Isolation",
     "Country outage does not affect other regions"),
    ("Auditability",
     "Queue depth reveals delivery delays per segment"),
    ("Durability",
     "durable:true + delivery_mode:2 — no message loss"),
]
add_textbox(slide6, "Operational Benefits", 7.2, 2.58, 5, 0.3,
            font_size=11, bold=True, color=NAVY)
for i, (b, d) in enumerate(benefits):
    y = 2.88 + i * 0.7
    add_rect(slide6, 7.2, y, 5.75, 0.62, LIGHT_GRAY)
    add_rect(slide6, 7.2, y, 0.05, 0.62, TEAL)
    add_textbox(slide6, b, 7.32, y + 0.04, 5.5, 0.25,
                font_size=10, bold=True, color=NAVY)
    add_textbox(slide6, d, 7.32, y + 0.3, 5.5, 0.25,
                font_size=9, color=MID_GRAY)

# Footer
add_rect(slide6, 0.35, 6.55, 12.6, 0.55, LIGHT_TEAL)
add_textbox(slide6,
            "Dead-letter queues capture all undeliverable messages — no silent failures.",
            0.5, 6.65, 12.2, 0.35,
            font_size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – RABBITMQ & QUEUE PROCESSORS
# ══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide7, WHITE)
header_bar(slide7, "RabbitMQ & Queue Processors",
           "Durable message bus with per-queue worker services")

# Left — RabbitMQ
add_textbox(slide7, "RabbitMQ Message Broker", 0.35, 1.55, 6, 0.35,
            font_size=13, bold=True, color=NAVY)
rb_items = [
    ("Persistent storage", "durable:true queues + delivery_mode:2 messages"),
    ("Management UI/API",  "Real-time queue depth & throughput monitoring"),
    ("AMQP protocol",      "Standard wire protocol for .NET consumers"),
    ("Dead-letter routing","Failed messages auto-routed to DLQ"),
    ("Horizontal scale",   "Clustering supported for HA deployments"),
]
for i, (lbl, val) in enumerate(rb_items):
    y = 1.98 + i * 0.68
    add_rect(slide7, 0.35, y, 6, 0.58, LIGHT_GRAY)
    add_rect(slide7, 0.35, y, 0.05, 0.58, TEAL)
    add_textbox(slide7, lbl, 0.45, y + 0.04, 2.4, 0.25,
                font_size=9.5, bold=True, color=NAVY)
    add_textbox(slide7, val, 0.45, y + 0.28, 5.8, 0.22,
                font_size=9, color=MID_GRAY)

# Right — Queue Processors
add_textbox(slide7, "Queue Processor Workers (×N)", 6.7, 1.55, 6.3, 0.35,
            font_size=13, bold=True, color=NAVY)
proc_items = [
    ("One per queue",         "Dedicated .NET Worker Service container"),
    ("Long-lived consumer",   "Single AMQP subscription — no polling overhead"),
    ("HTTP 2xx → ACK",        "Message permanently removed from queue"),
    ("HTTP 4xx → NACK",       "requeue=false → dead-letter queue"),
    ("HTTP 5xx → NACK",       "requeue=true → retry up to limit → DLQ"),
    ("Horizontal scaling",    "Multiple replicas share same queue (round-robin)"),
]
for i, (lbl, val) in enumerate(proc_items):
    y = 1.98 + i * 0.68
    add_rect(slide7, 6.7, y, 6.3, 0.58, LIGHT_GRAY)
    add_rect(slide7, 6.7, y, 0.05, 0.58, GREEN if "2xx" in lbl
             else (RED if "4xx" in lbl else (AMBER if "5xx" in lbl else TEAL)))
    add_textbox(slide7, lbl, 6.82, y + 0.04, 2.4, 0.25,
                font_size=9.5, bold=True, color=NAVY)
    add_textbox(slide7, val, 6.82, y + 0.28, 6.0, 0.22,
                font_size=9, color=MID_GRAY)

# Divider
add_rect(slide7, 6.55, 1.5, 0.04, 5.1, TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – SMS SPLITTER GATEWAY & MNO LAYER
# ══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide8, WHITE)
header_bar(slide8, "SMS Splitter Gateway & MNO Delivery Layer",
           "Intelligent routing from queue processor to mobile network operators")

# Splitter box
add_rect(slide8, 0.35, 1.55, 5.8, 4.55, LIGHT_GRAY)
add_rect(slide8, 0.35, 1.55, 0.05, 4.55, TEAL)
add_textbox(slide8, "SMS Splitter Gateway", 0.5, 1.62, 5.5, 0.32,
            font_size=13, bold=True, color=NAVY)
add_textbox(slide8, "C# .NET · IIS (Windows Server)", 0.5, 1.9, 5.5, 0.25,
            font_size=9, color=TEAL, italic=True)

splitter_points = [
    "Receives HTTP POST from Queue Processors",
    "Extracts destination MSISDN from payload",
    "Longest-prefix match: +25677 → Airtel Uganda",
    "Forwards to resolved MNO Passthrough API URL",
    "Centralises routing table — one place to update MNO endpoints",
    "Per-MNO throughput logging & latency metrics",
    "Optional fallback to secondary MNO on failure",
    "Decouples processors from MNO-specific details",
]
for i, pt in enumerate(splitter_points):
    add_textbox(slide8, f"• {pt}", 0.5, 2.25 + i * 0.46, 5.5, 0.38,
                font_size=9.5, color=DARK_GRAY)

# MNO Passthrough + SMSC box
add_rect(slide8, 6.5, 1.55, 6.5, 4.55, LIGHT_GRAY)
add_rect(slide8, 6.5, 1.55, 0.05, 4.55, NAVY)
add_textbox(slide8, "MNO Passthrough API + SMSC", 6.65, 1.62, 6.2, 0.32,
            font_size=13, bold=True, color=NAVY)
add_textbox(slide8, "Node.js · PM2 (Linux)  —  one instance per MNO",
            6.65, 1.9, 6.2, 0.25, font_size=9, color=TEAL, italic=True)

mno_sections = [
    ("Passthrough API", [
        "Accepts normalised HTTP SMS request",
        "Translates to MNO-specific protocol (SMPP / SDP)",
        "Submits to co-located SMSC module",
        "Returns delivery status to Splitter",
    ]),
    ("SMSC Module (SMPP)", [
        "Persistent TCP connection to MNO SMSC",
        "Manages bind/rebind, enquire_link keepalives",
        "Configurable window size & sequence tracking",
        "Exponential backoff rebind on session loss",
    ]),
    ("SMSC Module (SDP)", [
        "HTTP/HTTPS interface with auth & request signing",
        "Delivery receipts: correlates submit_sm_resp ACKs",
        "Processes deliver_sm receipts",
        "Supports SMPP 3.4 and SDP HTTP protocols",
    ]),
]
y = 2.25
for sec_title, pts in mno_sections:
    add_textbox(slide8, sec_title, 6.65, y, 6.2, 0.25,
                font_size=10, bold=True, color=TEAL)
    y += 0.28
    for pt in pts:
        add_textbox(slide8, f"• {pt}", 6.65, y, 6.2, 0.28,
                    font_size=9, color=DARK_GRAY)
        y += 0.28
    y += 0.1

# Arrow between boxes
add_rect(slide8, 6.18, 3.55, 0.28, 0.04, TEAL)
add_textbox(slide8, "▶", 6.38, 3.44, 0.25, 0.25, font_size=10, color=TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – DEPLOYMENT TOPOLOGY
# ══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide9, WHITE)
header_bar(slide9, "Deployment Topology",
           "Three-tier infrastructure — Windows IIS · Docker · Linux PM2")

tiers = [
    {
        "title": "Tier 1 — Windows Server (IIS)",
        "color": NAVY,
        "items": [
            "Elma Core (.NET API)",
            "SMS Splitter Gateway (.NET API)",
            "",
            "Role: Transaction processing & MNO routing",
            "Runtime: IIS on Windows Server",
            "Scaling: Vertical or load-balanced",
        ],
    },
    {
        "title": "Tier 2 — Docker Host",
        "color": RGBColor(0x00, 0x5F, 0x73),
        "items": [
            "SMS API (C# container)",
            "RabbitMQ API Wrapper (C# container)",
            "Queue Processors A … N (Worker Svc containers)",
            "RabbitMQ Broker",
            "",
            "Runtime: Docker (containerised)",
            "Scaling: Horizontal container replicas",
        ],
    },
    {
        "title": "Tier 3 — Linux Server (PM2)",
        "color": RGBColor(0x1A, 0x53, 0x36),
        "items": [
            "MNO-A Passthrough + SMSC (Node.js)",
            "MNO-B Passthrough + SMSC (Node.js)",
            "MNO-N … (one process per MNO)",
            "",
            "Runtime: PM2 process manager on Linux",
            "Scaling: PM2 multi-instance per MNO",
        ],
    },
]

tier_w = 4.0
for i, tier in enumerate(tiers):
    x = 0.35 + i * (tier_w + 0.33)
    # Header
    add_rect(slide9, x, 1.55, tier_w, 0.6, tier["color"])
    add_textbox(slide9, tier["title"], x + 0.1, 1.62, tier_w - 0.15, 0.42,
                font_size=10.5, bold=True, color=WHITE)
    # Body
    add_rect(slide9, x, 2.18, tier_w, 4.0, LIGHT_GRAY)
    add_rect(slide9, x, 2.18, 0.05, 4.0, tier["color"])
    for j, item in enumerate(tier["items"]):
        if item:
            prefix = "  " if item.startswith("Runtime") or item.startswith("Scaling") or item.startswith("Role") else "▸ "
            fc = MID_GRAY if item.startswith("Runtime") or item.startswith("Scaling") or item.startswith("Role") else DARK_GRAY
            add_textbox(slide9, prefix + item, x + 0.12, 2.28 + j * 0.5,
                        tier_w - 0.2, 0.4, font_size=9.5, color=fc)

# Arrows between tiers
add_textbox(slide9, "⟵  API calls  ⟶", 4.4, 4.0, 0.5, 0.35,
            font_size=8.5, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide9, "⟵  HTTP POST  ⟶", 8.72, 4.0, 0.5, 0.35,
            font_size=8.5, color=TEAL, align=PP_ALIGN.CENTER)

# Inter-tier dependency note
add_rect(slide9, 0.35, 6.38, 12.6, 0.7, NAVY)
add_textbox(slide9,
            "All tiers are independently deployable. "
            "Tier 1 (IIS) handles business logic · Tier 2 (Docker) handles async messaging · "
            "Tier 3 (Linux/PM2) handles MNO protocol integration.",
            0.5, 6.5, 12.2, 0.5,
            font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – FAILURE MODES & RESILIENCE
# ══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide10, WHITE)
header_bar(slide10, "Failure Modes & Resilience",
           "Every failure scenario has a defined mitigation strategy")

failures = [
    ("Elma Core unavailable",     "Transactions fail; no new SMS",         "Auto-restart + health-check",          RED),
    ("Queue API Wrapper down",    "Elma Core error; may retry internally",  "Container restart + Elma backoff",     AMBER),
    ("RabbitMQ broker down",      "No new messages; existing safe on disk", "Clustering + persistent disk storage", AMBER),
    ("Queue Processor crash",     "Messages unACKed, auto-requeued by broker","Auto-restart + consumer_timeout",   AMBER),
    ("SMS Splitter Gateway down", "5xx from processor → NACK & requeue",   "Backoff + dead-letter queue",          AMBER),
    ("MNO Passthrough down",      "Delivery fails; propagates to processor","PM2 auto-restart + fallback MNO",     AMBER),
    ("SMPP session lost",         "Delivery pauses until rebind",           "Exponential backoff rebind",           GREEN),
    ("MNO network outage",        "SMSC-level failure",                     "Dead-letter queue + scheduled retry",  RED),
]

# Column headers
cols_x = [0.35, 3.7, 7.3, 10.9]
col_labels = ["Failure", "Impact", "Mitigation", "Severity"]
col_widths = [3.25, 3.5, 3.5, 1.9]
add_rect(slide10, 0.35, 1.55, 12.6, 0.42, NAVY)
for lbl, x in zip(col_labels, cols_x):
    add_textbox(slide10, lbl, x + 0.08, 1.6, 3.2, 0.3,
                font_size=10, bold=True, color=WHITE)

for i, (failure, impact, mitigation, sev_color) in enumerate(failures):
    y = 2.02 + i * 0.58
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide10, 0.35, y, 12.6, 0.53, bg)
    add_textbox(slide10, failure,    0.43, y + 0.07, 3.2, 0.38, font_size=9.5, bold=True,  color=NAVY)
    add_textbox(slide10, impact,     3.78, y + 0.07, 3.4, 0.38, font_size=9,   color=MID_GRAY)
    add_textbox(slide10, mitigation, 7.38, y + 0.07, 3.4, 0.38, font_size=9,   color=DARK_GRAY)
    add_rect(slide10, 10.98, y + 0.12, 1.7, 0.3, sev_color)
    sev_lbl = "High" if sev_color == RED else ("Med" if sev_color == AMBER else "Low")
    add_textbox(slide10, sev_lbl, 10.98, y + 0.12, 1.7, 0.3,
                font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Legend
add_textbox(slide10, "■ High risk   ■ Medium risk   ■ Low / Auto-recovered",
            0.35, 6.7, 7, 0.35, font_size=9.5, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – NON-FUNCTIONAL CHARACTERISTICS
# ══════════════════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide11, WHITE)
header_bar(slide11, "Non-Functional Characteristics",
           "Throughput · Durability · Observability · Security")

nf_cards = [
    ("Throughput", TEAL, [
        "Queue Processor: bounded by SMS Splitter response time + prefetch count",
        "SMS Splitter Gateway: stateless HTTP handler — vertical scale or load balancer",
        "MNO Passthrough: bounded by SMPP window size (10–100 PDUs)",
        "PM2 supports multiple instances per MNO for higher throughput",
        "OTP queues: high prefetch  ·  PROMO queues: low prefetch",
    ]),
    ("Durability", GREEN, [
        "All RabbitMQ queues declared durable:true",
        "All messages sent with delivery_mode:2 (persisted to disk)",
        "Dead-letter queues prevent silent message loss",
        "SMPP persistent TCP with enquire_link keepalives",
        "Exponential backoff ensures no storm on reconnect",
    ]),
    ("Observability", NAVY, [
        "Queue depth via RabbitMQ Management UI / API",
        "Delivery latency: timestamped at enqueue & SMSC ACK",
        "SMPP session health via enquire_link cycle monitoring",
        "Per-MNO throughput logging in SMS Splitter Gateway",
        "Dead-letter queue depth = unresolved delivery failures",
    ]),
    ("Security & Ops", RGBColor(0x5B, 0x00, 0xB0), [
        "TLS on all inter-service HTTP connections",
        "SMPP bind credentials per MNO (not shared)",
        "SDP HTTP auth & request signing per MNO",
        "Container isolation for queue processors",
        "IIS on Windows Server for Elma Core (enterprise hardened)",
    ]),
]

card_w = 2.95
for i, (title, color, bullets) in enumerate(nf_cards):
    x = 0.35 + i * (card_w + 0.38)
    add_rect(slide11, x, 1.55, card_w, 0.5, color)
    add_textbox(slide11, title, x + 0.1, 1.62, card_w - 0.15, 0.35,
                font_size=12, bold=True, color=WHITE)
    add_rect(slide11, x, 2.08, card_w, 4.5, LIGHT_GRAY)
    add_rect(slide11, x, 2.08, 0.05, 4.5, color)
    for j, b in enumerate(bullets):
        add_textbox(slide11, f"• {b}", x + 0.12, 2.18 + j * 0.82, card_w - 0.2, 0.75,
                    font_size=9, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – KEY DESIGN DECISIONS
# ══════════════════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide12, WHITE)
header_bar(slide12, "Key Design Decisions",
           "Architecture choices and the problems they solve")

decisions = [
    ("Async-First Delivery",
     "Problem: MNO latency variability (50ms – 10s) would block transaction APIs.\n"
     "Decision: Enqueue SMS immediately; return 202 Accepted to caller.\n"
     "Benefit: Transaction layer never waits for carrier."),
    ("Queue-per-Segment",
     "Problem: A slow MNO in one country was starving OTP delivery for another bank.\n"
     "Decision: Isolate by COUNTRY × BANK_ID × CATEGORY.\n"
     "Benefit: Fault containment; independent priority tuning."),
    ("Protocol Abstraction (Splitter + Passthrough)",
     "Problem: Queue processors would need MNO-specific SMPP/SDP logic.\n"
     "Decision: Two-layer gateway (Splitter routes, Passthrough adapts).\n"
     "Benefit: New MNO integration = one new Passthrough instance, no processor changes."),
    ("Node.js for SMPP/SDP",
     "Problem: .NET SMPP libraries lacked mature async keep-alive handling for African MNOs.\n"
     "Decision: Node.js with event-loop for SMPP TCP persistence.\n"
     "Benefit: Non-blocking SMPP window management; PM2 restart on crash."),
]

for i, (title, body) in enumerate(decisions):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.45
    y = 1.6 + row * 2.6
    add_rect(slide12, x, y, 6.15, 2.4, LIGHT_GRAY)
    add_rect(slide12, x, y, 0.05, 2.4, TEAL)
    add_textbox(slide12, f"{i+1}. {title}", x + 0.12, y + 0.08, 5.9, 0.35,
                font_size=11, bold=True, color=NAVY)
    for j, line in enumerate(body.split("\n")):
        color = TEAL if line.startswith("Problem") else (DARK_GRAY if line.startswith("Decision") else GREEN)
        add_textbox(slide12, line, x + 0.12, y + 0.48 + j * 0.52, 5.9, 0.45,
                    font_size=9.5, color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – GLOSSARY
# ══════════════════════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide13, WHITE)
header_bar(slide13, "Glossary of Key Terms", "")

terms = [
    ("AMQP",         "Advanced Message Queuing Protocol — RabbitMQ wire protocol"),
    ("SMPP",         "Short Message Peer-to-Peer — binary protocol for SMSC submission"),
    ("SMSC",         "Short Message Service Centre — MNO system for SMS storage/forwarding"),
    ("SDP",          "Service Delivery Platform — HTTP/HTTPS SMS gateway"),
    ("MNO",          "Mobile Network Operator — telecom carrier (MTN, Airtel, Safaricom …)"),
    ("MSISDN",       "International mobile phone number (E.164 format)"),
    ("OTP",          "One-Time Password — time-sensitive auth SMS category"),
    ("PM2",          "Process Manager 2 for Node.js — keeps processes running, auto-restarts"),
    ("IIS",          "Internet Information Services — Microsoft enterprise web server"),
    ("DLQ",          "Dead-letter queue — holds failed messages for manual inspection/retry"),
    ("Worker Service","Background .NET long-running service — no HTTP listener, just consumers"),
    ("Prefetch count","Number of unACKed messages a consumer holds — RabbitMQ flow control"),
    ("Window size",   "SMPP: max outstanding PDUs before awaiting acknowledgement"),
    ("NACK",          "Negative acknowledgement — tells broker message was not processed"),
]

col1 = terms[:7]
col2 = terms[7:]
for i, (term, defn) in enumerate(col1):
    y = 1.6 + i * 0.68
    add_rect(slide13, 0.35, y, 6.1, 0.6, LIGHT_GRAY)
    add_rect(slide13, 0.35, y, 0.05, 0.6, TEAL)
    add_textbox(slide13, term, 0.47, y + 0.04, 1.5, 0.25,
                font_size=10, bold=True, color=NAVY)
    add_textbox(slide13, defn, 0.47, y + 0.28, 5.8, 0.25,
                font_size=9, color=MID_GRAY)
for i, (term, defn) in enumerate(col2):
    y = 1.6 + i * 0.68
    add_rect(slide13, 6.85, y, 6.1, 0.6, LIGHT_GRAY)
    add_rect(slide13, 6.85, y, 0.05, 0.6, TEAL)
    add_textbox(slide13, term, 6.97, y + 0.04, 1.6, 0.25,
                font_size=10, bold=True, color=NAVY)
    add_textbox(slide13, defn, 6.97, y + 0.28, 5.8, 0.25,
                font_size=9, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CLOSING / THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(BLANK)
fill_slide_bg(slide14, NAVY)

add_rect(slide14, 0, 5.5, 13.33, 0.15, TEAL)

add_textbox(slide14, "ELMA", 0.5, 0.5, 4, 0.7,
            font_size=40, bold=True, color=TEAL)
add_textbox(slide14, "SMS Platform Architecture", 0.5, 1.2, 12, 0.65,
            font_size=30, bold=True, color=WHITE)
add_textbox(slide14,
            "Multi-Channel  ·  Queue-Isolated  ·  Protocol-Agnostic  ·  Resilient",
            0.5, 1.85, 12, 0.45,
            font_size=15, italic=True, color=LIGHT_TEAL)

summary_points = [
    "Three entry channels (App, USSD, Web) converge on one async pipeline",
    "Granular queue isolation per Country × Bank × Category",
    "Async-first: MNO back-pressure never impacts transaction latency",
    "Dead-letter queues guarantee zero silent message loss",
    "Three-tier deployment: IIS · Docker · Linux/PM2 — independently scalable",
]
for i, pt in enumerate(summary_points):
    add_textbox(slide14, f"✓  {pt}", 0.5, 2.6 + i * 0.52, 12, 0.42,
                font_size=11.5, color=WHITE)

add_textbox(slide14, "Questions & Discussion", 0.5, 5.65, 12, 0.5,
            font_size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(slide14, "June 2026  ·  StayHere Engineering",
            0.5, 6.2, 12, 0.35,
            font_size=10, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
output_path = r"d:\samples\AG\repos\StayHereMVP\docs\SMS Architecture\Elma_SMS_Architecture_Board_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
